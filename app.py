from flask import Flask, render_template, request, jsonify
import os
import logging
from datetime import datetime
import traceback
import time
from app.utils.converters import (
    convert_docx, 
    convert_xlsx, 
    convert_pptx, 
    convert_pdf, 
    convert_txt, 
    convert_md,
    convert_to_markdown
)
from app.utils.converter_factory import converter_factory
from app.utils.url_converter import URLConverter
# 导入flasgger相关模块
from flasgger import Swagger, swag_from
import re

# 配置日志
log_dir = 'logs'
os.makedirs(log_dir, exist_ok=True)
log_file = os.path.join(log_dir, f'app_{datetime.now().strftime("%Y%m%d")}.log')

logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler(log_file, encoding='utf-8'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

app = Flask(__name__, template_folder='app/templates', static_folder='app/static')
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 限制上传文件大小为50MB
app.config['UPLOAD_FOLDER'] = 'uploads'

# 配置Swagger
swagger_config = {
    "headers": [],
    "specs": [
        {
            "endpoint": "apispec",
            "route": "/apispec.json",
            "rule_filter": lambda rule: True,
            "model_filter": lambda tag: True,
        }
    ],
    "static_url_path": "/flasgger_static",
    "swagger_ui": True,
    "specs_route": "/swagger/",
    "uiversion": 3,
    "ui_params": {
        "defaultModelsExpandDepth": -1,  # 不显示Models部分
        "docExpansion": "list", # 默认展开列表
        "displayRequestDuration": True, # 显示请求持续时间
        "defaultModelRendering": "model" # 使用模型渲染
    }
}

# 添加国际化支持
swagger_template = {
    "swagger": "2.0",
    "info": {
        "title": "X2Knowledge API",
        "description": "用于文档转换的API接口",
        "version": "1.0.0",
        "contact": {
            "name": "X2Knowledge API",
        }
    },
    "x-i18n": {
        "zh": {
            "info": {
                "title": "X2Knowledge API",
                "description": "用于文档转换的API接口"
            }
        },
        "en": {
            "info": {
                "title": "Document Conversion API",
                "description": "API interfaces for document conversion"
            }
        }
    }
}

swagger = Swagger(app, config=swagger_config, template=swagger_template)

# 确保上传文件夹存在
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

@app.route('/')
def index():
    logger.info("访问首页")
    return render_template('index.html')

@app.route('/OCR安装说明.md')
def ocr_installation():
    logger.info("访问OCR安装说明")
    try:
        with open('OCR安装说明.md', 'r', encoding='utf-8') as f:
            content = f.read()
        return content, 200, {'Content-Type': 'text/markdown'}
    except Exception as e:
        logger.error(f"读取OCR安装说明失败: {str(e)}")
        return "OCR安装说明文件不存在", 404

@app.route('/api-docs')
def api_docs():
    logger.info("访问API文档页面")
    return render_template('api-docs.html')

# 转换为文本
@app.route('/convert', methods=['POST'])
def convert():
    if 'file' not in request.files:
        logger.warning("没有文件上传")
        return jsonify({'error': '没有文件上传'}), 400
    
    file = request.files['file']
    
    if file.filename == '':
        logger.warning("未选择文件")
        return jsonify({'error': '未选择文件'}), 400
    
    # 获取文件扩展名
    filename = file.filename
    file_ext = os.path.splitext(filename)[1].lower()
    
    logger.info(f"接收到文件: {filename}, 类型: {file_ext}")
    
    # 保存上传的文件
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    try:
        file.save(file_path)
        logger.info(f"文件保存成功: {file_path}")
        
        # 检查文件大小
        file_size = os.path.getsize(file_path)
        logger.info(f"文件大小: {file_size / 1024:.2f} KB")
        
        # 根据文件类型调用相应的转换函数
        if file_ext in ['.doc', '.docx']:
            logger.info(f"开始转换Word文档: {filename}")
            text = convert_docx(file_path)
        elif file_ext in ['.xls', '.xlsx']:
            logger.info(f"开始转换Excel文件: {filename}")
            text = convert_xlsx(file_path)
        elif file_ext in ['.ppt', '.pptx']:
            logger.info(f"开始转换PowerPoint文件: {filename}")
            text = convert_pptx(file_path)
        elif file_ext == '.pdf':
            logger.info(f"开始转换PDF文件: {filename}")
            text = convert_pdf(file_path)
        elif file_ext == '.txt':
            logger.info(f"开始读取文本文件: {filename}")
            text = convert_txt(file_path)
        elif file_ext == '.md':
            logger.info(f"开始转换Markdown文件: {filename}")
            text = convert_md(file_path)
        elif file_ext in ['.mp3', '.wav']:
            logger.info(f"开始转换音频文件: {filename}")
            text = "音频文件支持Markdown格式导出，请使用转MD功能"
        else:
            logger.warning(f"不支持的文件格式: {file_ext}")
            os.remove(file_path)  # 删除不支持的文件
            return jsonify({'error': f'不支持的文件格式: {file_ext}'}), 400
        
        # 记录转换结果
        text_length = len(text)
        logger.info(f"文件转换成功, 文本长度: {text_length} 字符")
        
        # 删除临时文件
        try:
            os.remove(file_path)
            logger.info(f"临时文件已删除: {file_path}")
        except Exception as e:
            logger.warning(f"无法删除临时文件: {file_path}, 原因: {str(e)}")
            # 尝试强制关闭文件句柄后再删除
            try:
                import gc
                gc.collect()  # 触发垃圾回收，释放未关闭的文件句柄
                os.remove(file_path)
                logger.info(f"临时文件已删除(第二次尝试): {file_path}")
            except Exception as e2:
                logger.error(f"无法删除临时文件(第二次尝试): {file_path}, 原因: {str(e2)}")
                # 不再尝试删除，避免阻塞程序
        
        return jsonify({'text': text})

    except Exception as e:
        # 详细记录异常信息
        error_msg = str(e)
        logger.error(f"转换失败: {error_msg}")
        logger.error(traceback.format_exc())
        
        # 出错时删除临时文件
        if os.path.exists(file_path):
            try:
                os.remove(file_path)
                logger.info(f"清理临时文件: {file_path}")
            except Exception as del_e:
                logger.warning(f"无法删除临时文件: {file_path}, 原因: {str(del_e)}")
                # 不再尝试删除，避免阻塞程序
            
        return jsonify({
            'error': f'转换失败: {error_msg}', 
            'details': traceback.format_exc()
        }), 500

# 转换为Markdown
@app.route('/convert-to-md', methods=['POST'])
def convert_md_route():
    if 'file' not in request.files:
        logger.warning("没有文件上传")
        return jsonify({'error': '没有文件上传'}), 400
    
    file = request.files['file']
    
    if file.filename == '':
        logger.warning("未选择文件")
        return jsonify({'error': '未选择文件'}), 400
    
    # 获取文件扩展名
    filename = file.filename
    file_ext = os.path.splitext(filename)[1].lower()
    
    logger.info(f"接收到文件(转MD): {filename}, 类型: {file_ext}")
    
    # 检查文件扩展名
    supported_extensions = ['.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx', '.pdf', '.txt', '.md', '.mp3', '.wav']
    if file_ext not in supported_extensions:
        logger.warning(f"不支持的文件格式: {file_ext}")
        return jsonify({'error': f'不支持的文件格式: {file_ext}'}), 400
    
    # 保存上传的文件
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    try:
        file.save(file_path)
        logger.info(f"文件保存成功(转MD): {file_path}")
        
        # 检查文件大小
        file_size = os.path.getsize(file_path)
        logger.info(f"文件大小(转MD): {file_size / 1024:.2f} KB")
        
        # 使用MarkItDown转换为Markdown
        logger.info(f"开始将文件转换为Markdown: {filename}")
        markdown_text = convert_to_markdown(file_path)
        
        # 记录转换结果
        text_length = len(markdown_text)
        logger.info(f"文件转换为Markdown成功, 文本长度: {text_length} 字符")
        
        # 删除临时文件
        try:
            os.remove(file_path)
            logger.info(f"临时文件已删除(转MD): {file_path}")
        except Exception as e:
            logger.warning(f"无法删除临时文件(转MD): {file_path}, 原因: {str(e)}")
            try:
                import gc
                gc.collect()
                os.remove(file_path)
                logger.info(f"临时文件已删除(转MD)(第二次尝试): {file_path}")
            except Exception as e2:
                logger.error(f"无法删除临时文件(转MD)(第二次尝试): {file_path}, 原因: {str(e2)}")
        
        return jsonify({'text': markdown_text})

    except Exception as e:
        # 详细记录异常信息
        error_msg = str(e)
        logger.error(f"转换为Markdown失败: {error_msg}")
        logger.error(traceback.format_exc())
        
        # 出错时删除临时文件
        if os.path.exists(file_path):
            try:
                os.remove(file_path)
                logger.info(f"清理临时文件(转MD): {file_path}")
            except Exception as del_e:
                logger.warning(f"无法删除临时文件(转MD): {file_path}, 原因: {str(del_e)}")
            
        return jsonify({
            'error': f'转换为Markdown失败: {error_msg}', 
            'details': traceback.format_exc()
        }), 500

# API: 基本文本转换
@app.route('/api/convert', methods=['POST'])
@swag_from('swagger_docs/convert.yml')
def api_convert():
    start_time = time.time()
    
    if 'file' not in request.files:
        logger.warning("API调用：没有文件上传")
        return jsonify({'error': '没有文件上传'}), 400
    
    file = request.files['file']
    
    if file.filename == '':
        logger.warning("API调用：未选择文件")
        return jsonify({'error': '未选择文件'}), 400
    
    # 获取文件扩展名
    filename = file.filename
    file_ext = os.path.splitext(filename)[1].lower()
    
    logger.info(f"API调用：接收到文件: {filename}, 类型: {file_ext}")
    
    # 保存上传的文件
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    try:
        file.save(file_path)
        logger.info(f"API调用：文件保存成功: {file_path}")
        
        # 检查文件大小
        file_size = os.path.getsize(file_path)
        logger.info(f"API调用：文件大小: {file_size / 1024:.2f} KB")
        
        # 根据文件类型调用相应的转换函数
        if file_ext in ['.doc', '.docx']:
            logger.info(f"API调用：开始转换Word文档: {filename}")
            text = convert_docx(file_path)
        elif file_ext in ['.xls', '.xlsx']:
            logger.info(f"API调用：开始转换Excel文件: {filename}")
            text = convert_xlsx(file_path)
        elif file_ext in ['.ppt', '.pptx']:
            logger.info(f"API调用：开始转换PowerPoint文件: {filename}")
            text = convert_pptx(file_path)
        elif file_ext == '.pdf':
            logger.info(f"API调用：开始转换PDF文件: {filename}")
            text = convert_pdf(file_path)
        elif file_ext == '.txt':
            logger.info(f"API调用：开始转换文本文件: {filename}")
            text = convert_txt(file_path)
        elif file_ext == '.md':
            logger.info(f"API调用：开始转换Markdown文件: {filename}")
            text = convert_md(file_path)
        elif file_ext in ['.mp3', '.wav']:
            logger.info(f"API调用：开始转换音频文件: {filename}")
            text = "音频文件支持Markdown格式导出，请使用转MD API"
        else:
            error_msg = f"不支持的文件类型: {file_ext}"
            logger.warning(f"API调用：{error_msg}")
            
            # 删除临时文件
            try:
                os.remove(file_path)
                logger.info(f"API调用：临时文件已删除: {file_path}")
            except Exception as e:
                logger.warning(f"API调用：无法删除临时文件: {file_path}, 原因: {str(e)}")
            
            return jsonify({'error': error_msg}), 400
        
        processing_time = time.time() - start_time
        logger.info(f"API调用：文件转换完成，耗时: {processing_time:.2f}秒")
        
        # 删除临时文件
        try:
            os.remove(file_path)
            logger.info(f"API调用：临时文件已删除: {file_path}")
        except Exception as e:
            logger.warning(f"API调用：无法删除临时文件: {file_path}, 原因: {str(e)}")
            # 尝试强制关闭文件句柄后再删除
            try:
                import gc
                gc.collect()  # 触发垃圾回收，释放未关闭的文件句柄
                os.remove(file_path)
                logger.info(f"API调用：临时文件已删除(第二次尝试): {file_path}")
            except Exception as e2:
                logger.error(f"API调用：无法删除临时文件(第二次尝试): {file_path}, 原因: {str(e2)}")
                # 不再尝试删除，避免阻塞程序
        
        # 返回API响应
        return jsonify({
            'text': text,
            'filename': filename,
            'file_size': file_size,
            'processing_time': round(processing_time, 2)
        })
    except Exception as e:
        error_msg = str(e)
        logger.error(f"API调用：转换失败: {error_msg}")
        logger.error(traceback.format_exc())
        
        # 出错时删除临时文件
        if os.path.exists(file_path):
            try:
                os.remove(file_path)
                logger.info(f"API调用：清理临时文件: {file_path}")
            except Exception as del_e:
                logger.warning(f"API调用：无法删除临时文件: {file_path}, 原因: {str(del_e)}")
        
        return jsonify({
            'error': error_msg,
            'details': traceback.format_exc()
        }), 500

# API：转换为Markdown
@app.route('/api/convert-to-md', methods=['POST'])
@swag_from('swagger_docs/convert_to_md.yml')
def api_convert_md():
    start_time = time.time()
    
    if 'file' not in request.files:
        logger.warning("API调用(转MD)：没有文件上传")
        return jsonify({'error': '没有文件上传'}), 400
    
    file = request.files['file']
    
    if file.filename == '':
        logger.warning("API调用(转MD)：未选择文件")
        return jsonify({'error': '未选择文件'}), 400
    
    # 获取文件扩展名
    filename = file.filename
    file_ext = os.path.splitext(filename)[1].lower()
    
    logger.info(f"API调用(转MD)：接收到文件: {filename}, 类型: {file_ext}")
    
    # 检查文件扩展名
    supported_extensions = ['.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx', '.pdf', '.txt', '.md', '.mp3', '.wav']
    if file_ext not in supported_extensions:
        logger.warning(f"API调用(转MD)：不支持的文件格式: {file_ext}")
        return jsonify({'error': f'不支持的文件格式: {file_ext}'}), 400
    
    # 保存上传的文件
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    try:
        file.save(file_path)
        logger.info(f"API调用(转MD)：文件保存成功: {file_path}")
        
        # 检查文件大小
        file_size = os.path.getsize(file_path)
        logger.info(f"API调用(转MD)：文件大小: {file_size / 1024:.2f} KB")
        
        # 使用MarkItDown转换为Markdown
        logger.info(f"API调用(转MD)：开始将文件转换为Markdown: {filename}")
        markdown_text = convert_to_markdown(file_path)
        
        processing_time = time.time() - start_time
        logger.info(f"API调用(转MD)：文件转换完成，耗时: {processing_time:.2f}秒")
        
        # 删除临时文件
        try:
            os.remove(file_path)
            logger.info(f"API调用(转MD)：临时文件已删除: {file_path}")
        except Exception as e:
            logger.warning(f"API调用(转MD)：无法删除临时文件: {file_path}, 原因: {str(e)}")
            try:
                import gc
                gc.collect()
                os.remove(file_path)
                logger.info(f"API调用(转MD)：临时文件已删除(第二次尝试): {file_path}")
            except Exception as e2:
                logger.error(f"API调用(转MD)：无法删除临时文件(第二次尝试): {file_path}, 原因: {str(e2)}")
        
        # 返回API响应
        return jsonify({
            'text': markdown_text,
            'filename': filename,
            'file_size': file_size,
            'processing_time': round(processing_time, 2),
            'converter': 'docling'
        })
    except Exception as e:
        error_msg = str(e)
        logger.error(f"API调用(转MD)：转换失败: {error_msg}")
        logger.error(traceback.format_exc())
        
        # 出错时删除临时文件
        if os.path.exists(file_path):
            try:
                os.remove(file_path)
                logger.info(f"API调用(转MD)：清理临时文件: {file_path}")
            except Exception as del_e:
                logger.warning(f"API调用(转MD)：无法删除临时文件: {file_path}, 原因: {str(del_e)}")
        
        return jsonify({
            'error': error_msg,
            'details': traceback.format_exc()
        }), 500

@app.route('/api/convert-to-md-images-file-docling', methods=['POST'])
@swag_from('swagger_docs/convert_to_md_images_file_docling.yml')
def api_convert_to_md_images_file_docling():
    """
    使用Docling将文件转换为Markdown并导出图片
    """
    start_time = time.time()
    execution_id = datetime.now().strftime("%Y%m%d%H%M%S")
    
    # 获取上传的文件
    if 'file' not in request.files:
        logger.warning(f"API调用(Docling图片导出){execution_id}：没有文件上传")
        return jsonify({'error': '没有文件上传'}), 400
    
    file = request.files['file']
    
    if file.filename == '':
        logger.warning(f"API调用(Docling图片导出){execution_id}：未选择文件")
        return jsonify({'error': '未选择文件'}), 400
    
    # 获取输出目录
    output_dir = request.form.get('output_dir')
    if not output_dir:
        logger.warning(f"API调用(Docling图片导出){execution_id}：未指定输出目录")
        return jsonify({'error': '未指定输出目录'}), 400
    
    # 获取文件扩展名
    filename = file.filename
    base_filename = os.path.splitext(filename)[0]
    file_ext = os.path.splitext(filename)[1].lower()
    
    logger.info(f"API调用(Docling图片导出){execution_id}：接收到文件: {filename}, 类型: {file_ext}")
    
    # 初始化Docling转换器
    docling_converter = converter_factory.get_converter("docling")
    
    # 检查文件扩展名是否支持
    if not docling_converter.is_format_supported(file_ext):
        error_msg = f"Docling不支持的文件类型: {file_ext}"
        logger.warning(f"API调用(Docling图片导出){execution_id}：{error_msg}")
        return jsonify({'error': error_msg}), 400
    
    # 保存上传的文件
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    try:
        file.save(file_path)
        logger.info(f"API调用(Docling图片导出){execution_id}：文件保存成功: {file_path}")
        
        # 检查文件大小
        file_size = os.path.getsize(file_path)
        logger.info(f"API调用(Docling图片导出){execution_id}：文件大小: {file_size / 1024:.2f} KB")
        
        # 检查输出目录是否存在，不存在则创建
        if not os.path.exists(output_dir):
            try:
                os.makedirs(output_dir, exist_ok=True)
                logger.info(f"API调用(Docling图片导出){execution_id}：创建输出目录: {output_dir}")
            except Exception as e:
                error_msg = f"无法创建输出目录: {str(e)}"
                logger.error(f"API调用(Docling图片导出){execution_id}：{error_msg}")
                
                # 删除临时文件
                try:
                    os.remove(file_path)
                    logger.info(f"API调用(Docling图片导出){execution_id}：临时文件已删除: {file_path}")
                except Exception as del_e:
                    logger.warning(f"API调用(Docling图片导出){execution_id}：无法删除临时文件: {file_path}, 原因: {str(del_e)}")
                
                return jsonify({'error': error_msg}), 500
        
        # 使用Docling处理文件
        try:
            from docling.document_converter import DocumentConverter
            from docling.datamodel.pipeline_options import PdfPipelineOptions
            from docling.datamodel.base_models import InputFormat, FigureElement
            from docling.document_converter import PdfFormatOption
            from docling_core.types.doc import ImageRefMode, PictureItem, TableItem
            
            # 设置Docling选项，启用图片导出
            IMAGE_RESOLUTION_SCALE = 2.0
            pipeline_options = PdfPipelineOptions()
            pipeline_options.images_scale = IMAGE_RESOLUTION_SCALE
            pipeline_options.generate_page_images = True
            pipeline_options.generate_picture_images = True
            
            # 创建通用格式选项
            format_options = {}
            
            # 根据文件类型添加适当的格式选项
            if file_ext.lower() in ['.pdf']:
                format_options[InputFormat.PDF] = PdfFormatOption(pipeline_options=pipeline_options)
            # 对于其他格式，不设置特定选项，使用Docling默认设置
            
            doc_converter = DocumentConverter(format_options=format_options)
            
            logger.info(f"API调用(Docling图片导出){execution_id}：开始转换文件并导出图片")
            doc_converter = DocumentConverter(
                format_options={
                    InputFormat.PDF: PdfFormatOption(pipeline_options=pipeline_options)
                }
            )
            
            logger.info(f"API调用(Docling图片导出){execution_id}：开始转换文件并导出图片")
            
            # 转换文件
            conv_res = doc_converter.convert(file_path)
            doc_filename = conv_res.input.file.stem
            
            # 保存Markdown文件
            md_filename = os.path.join(output_dir, f"{doc_filename}.md")
            conv_res.document.save_as_markdown(md_filename, image_mode=ImageRefMode.REFERENCED)
            logger.info(f"API调用(Docling图片导出){execution_id}：Markdown已保存到: {md_filename}")
            
            # 保存页面图片
            page_image_paths = []
            for page_no, page in conv_res.document.pages.items():
                page_no = page.page_no
                page_image_filename = os.path.join(output_dir, f"{doc_filename}-page-{page_no}.png")
                try:
                    if hasattr(page, 'image') and page.image is not None and hasattr(page.image, 'pil_image') and page.image.pil_image is not None:
                        with open(page_image_filename, 'wb') as fp:
                            page.image.pil_image.save(fp, format="PNG")
                        page_image_paths.append(page_image_filename)
                        logger.info(f"API调用(Docling图片导出){execution_id}：页面图片已保存: {page_image_filename}")
                    else:
                        logger.warning(f"API调用(Docling图片导出){execution_id}：页面 {page_no} 没有可用的图片")
                except Exception as e:
                    logger.warning(f"API调用(Docling图片导出){execution_id}：保存页面 {page_no} 图片时出错: {str(e)}")
            
            # 保存图片和表格
            table_counter = 0  # 实际成功提取的表格数量
            picture_counter = 0  # 实际成功提取的图片数量
            table_image_paths = []
            picture_image_paths = []
            
            for element, _level in conv_res.document.iterate_items():
                if isinstance(element, TableItem):
                    try:
                        image = element.get_image(conv_res.document)
                        if image is not None:
                            table_counter += 1  # 只有当成功获取图片时才增加计数
                            element_image_filename = os.path.join(output_dir, f"{doc_filename}-table-{table_counter}.png")
                            with open(element_image_filename, 'wb') as fp:
                                image.save(fp, "PNG")
                            table_image_paths.append(element_image_filename)
                            logger.info(f"API调用(Docling图片导出){execution_id}：表格图片已保存: {element_image_filename}")
                        else:
                            logger.warning(f"API调用(Docling图片导出){execution_id}：无法获取表格图片，图片为None")
                    except Exception as e:
                        logger.warning(f"API调用(Docling图片导出){execution_id}：处理表格图片时出错: {str(e)}")
                
                if isinstance(element, PictureItem):
                    try:
                        image = element.get_image(conv_res.document)
                        if image is not None:
                            picture_counter += 1  # 只有当成功获取图片时才增加计数
                            element_image_filename = os.path.join(output_dir, f"{doc_filename}-picture-{picture_counter}.png")
                            with open(element_image_filename, 'wb') as fp:
                                image.save(fp, "PNG")
                            picture_image_paths.append(element_image_filename)
                            logger.info(f"API调用(Docling图片导出){execution_id}：图片已保存: {element_image_filename}")
                        else:
                            logger.warning(f"API调用(Docling图片导出){execution_id}：无法获取图片，图片为None")
                    except Exception as e:
                        logger.warning(f"API调用(Docling图片导出){execution_id}：处理图片时出错: {str(e)}")
            
            # 读取Markdown内容
            with open(md_filename, 'r', encoding='utf-8') as f:
                markdown_text = f.read()
            
            processing_time = time.time() - start_time
            logger.info(f"API调用(Docling图片导出){execution_id}：处理完成，耗时: {processing_time:.2f}秒")
            
            # 返回处理结果
            return jsonify({
                'output_path': md_filename,
                'filename': filename,
                'file_size': file_size,
                'processing_time': round(processing_time, 2),
                'converter': 'docling',
                'page_count': len(page_image_paths),
                'table_count': table_counter,
                'picture_count': picture_counter,
                'page_images': page_image_paths,
                'table_images': table_image_paths,
                'picture_images': picture_image_paths
            })
            
        except ImportError as e:
            error_msg = f"Docling库导入错误: {str(e)}"
            logger.error(f"API调用(Docling图片导出){execution_id}：{error_msg}")
            return jsonify({'error': error_msg}), 500
        except Exception as e:
            error_msg = f"Docling处理失败: {str(e)}"
            logger.error(f"API调用(Docling图片导出){execution_id}：{error_msg}")
            logger.error(traceback.format_exc())
            return jsonify({'error': error_msg, 'details': traceback.format_exc()}), 500
    
    except Exception as e:
        error_msg = f"处理文件时出错: {str(e)}"
        logger.error(f"API调用(Docling图片导出){execution_id}：{error_msg}")
        logger.error(traceback.format_exc())
        
        # 删除临时文件
        if os.path.exists(file_path):
            try:
                os.remove(file_path)
                logger.info(f"API调用(Docling图片导出){execution_id}：临时文件已删除: {file_path}")
            except Exception as del_e:
                logger.warning(f"API调用(Docling图片导出){execution_id}：无法删除临时文件: {file_path}, 原因: {str(del_e)}")
        
        return jsonify({'error': error_msg, 'details': traceback.format_exc()}), 500
    
    finally:
        # 删除临时文件
        if os.path.exists(file_path):
            try:
                os.remove(file_path)
                logger.info(f"API调用(Docling图片导出){execution_id}：临时文件已删除: {file_path}")
            except Exception as del_e:
                logger.warning(f"API调用(Docling图片导出){execution_id}：无法删除临时文件: {file_path}, 原因: {str(del_e)}")

# Docling转换为Markdown
@app.route('/convert-to-md-docling', methods=['POST'])
def convert_md_docling_route():
    if 'file' not in request.files:
        logger.warning("没有文件上传(Docling)")
        return jsonify({'error': '没有文件上传'}), 400
    
    file = request.files['file']
    
    if file.filename == '':
        logger.warning("未选择文件(Docling)")
        return jsonify({'error': '未选择文件'}), 400
    
    # 获取文件扩展名
    filename = file.filename
    file_ext = os.path.splitext(filename)[1].lower()
    
    logger.info(f"接收到文件(Docling转MD): {filename}, 类型: {file_ext}")
    
    # 获取docling转换器
    docling_converter = converter_factory.get_converter('docling')
    if not hasattr(docling_converter, 'is_available') or not docling_converter.is_available:
        logger.error("Docling转换器不可用")
        return jsonify({'error': 'Docling转换器不可用，请确认已安装docling库'}), 500
    
    # 检查文件扩展名是否支持
    if not docling_converter.is_format_supported(file_ext):
        logger.warning(f"Docling不支持的文件格式: {file_ext}")
        return jsonify({'error': f'Docling不支持的文件格式: {file_ext}'}), 400
    
    # 保存上传的文件
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    try:
        file.save(file_path)
        logger.info(f"文件保存成功(Docling转MD): {file_path}")
        
        # 检查文件大小
        file_size = os.path.getsize(file_path)
        logger.info(f"文件大小(Docling转MD): {file_size / 1024:.2f} KB")
        
        # 使用Docling转换为Markdown
        logger.info(f"开始使用Docling将文件转换为Markdown: {filename}")
        markdown_text = docling_converter.convert_to_markdown(file_path)
        
        # 记录转换结果
        text_length = len(markdown_text)
        logger.info(f"文件使用Docling转换为Markdown成功, 文本长度: {text_length} 字符")
        
        # 删除临时文件
        try:
            os.remove(file_path)
            logger.info(f"临时文件已删除(Docling转MD): {file_path}")
        except Exception as e:
            logger.warning(f"无法删除临时文件(Docling转MD): {file_path}, 原因: {str(e)}")
            try:
                import gc
                gc.collect()
                os.remove(file_path)
                logger.info(f"临时文件已删除(Docling转MD)(第二次尝试): {file_path}")
            except Exception as e2:
                logger.error(f"无法删除临时文件(Docling转MD)(第二次尝试): {file_path}, 原因: {str(e2)}")
        
        return jsonify({'text': markdown_text})

    except Exception as e:
        # 详细记录异常信息
        error_msg = str(e)
        logger.error(f"Docling转换为Markdown失败: {error_msg}")
        logger.error(traceback.format_exc())
        
        # 出错时删除临时文件
        if os.path.exists(file_path):
            try:
                os.remove(file_path)
                logger.info(f"清理临时文件(Docling转MD): {file_path}")
            except Exception as del_e:
                logger.warning(f"无法删除临时文件(Docling转MD): {file_path}, 原因: {str(del_e)}")
            
        return jsonify({
            'error': f'Docling转换为Markdown失败: {error_msg}', 
            'details': traceback.format_exc()
        }), 500

# API：使用Docling转换为Markdown
@app.route('/api/convert-to-md-docling', methods=['POST'])
@swag_from('swagger_docs/convert_to_md_docling.yml')
def api_convert_md_docling():
    start_time = time.time()
    
    if 'file' not in request.files:
        logger.warning("API调用(Docling)：没有文件上传")
        return jsonify({'error': '没有文件上传'}), 400
    
    file = request.files['file']
    
    if file.filename == '':
        logger.warning("API调用(Docling)：未选择文件")
        return jsonify({'error': '未选择文件'}), 400
    
    # 获取文件扩展名
    filename = file.filename
    file_ext = os.path.splitext(filename)[1].lower()
    
    logger.info(f"API调用(Docling)：接收到文件: {filename}, 类型: {file_ext}")
    
    # 获取docling转换器
    docling_converter = converter_factory.get_converter('docling')
    if not hasattr(docling_converter, 'is_available') or not docling_converter.is_available:
        logger.error("API调用(Docling)：Docling转换器不可用")
        return jsonify({'error': 'Docling转换器不可用，请确认已安装docling库'}), 500
    
    # 检查文件扩展名是否支持
    if not docling_converter.is_format_supported(file_ext):
        error_msg = f"Docling不支持的文件类型: {file_ext}"
        logger.warning(f"API调用(Docling)：{error_msg}")
        return jsonify({'error': error_msg}), 400
    
    # 保存上传的文件
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    try:
        file.save(file_path)
        logger.info(f"API调用(Docling)：文件保存成功: {file_path}")
        
        # 检查文件大小
        file_size = os.path.getsize(file_path)
        logger.info(f"API调用(Docling)：文件大小: {file_size / 1024:.2f} KB")
        
        # 使用Docling转换为Markdown
        logger.info(f"API调用(Docling)：开始转换文件: {filename}")
        markdown_text = docling_converter.convert_to_markdown(file_path)
        
        processing_time = time.time() - start_time
        logger.info(f"API调用(Docling)：文件转换完成，耗时: {processing_time:.2f}秒")
        
        # 删除临时文件
        try:
            os.remove(file_path)
            logger.info(f"API调用(Docling)：临时文件已删除: {file_path}")
        except Exception as e:
            logger.warning(f"API调用(Docling)：无法删除临时文件: {file_path}, 原因: {str(e)}")
            try:
                import gc
                gc.collect()
                os.remove(file_path)
                logger.info(f"API调用(Docling)：临时文件已删除(第二次尝试): {file_path}")
            except Exception as e2:
                logger.error(f"API调用(Docling)：无法删除临时文件(第二次尝试): {file_path}, 原因: {str(e2)}")
        
        # 返回API响应
        return jsonify({
            'text': markdown_text,
            'filename': filename,
            'file_size': file_size,
            'processing_time': round(processing_time, 2),
            'converter': 'docling'
        })
    except Exception as e:
        error_msg = str(e)
        logger.error(f"API调用(Docling)：转换失败: {error_msg}")
        logger.error(traceback.format_exc())
        
        # 出错时删除临时文件
        if os.path.exists(file_path):
            try:
                os.remove(file_path)
                logger.info(f"API调用(Docling)：清理临时文件: {file_path}")
            except Exception as del_e:
                logger.warning(f"API调用(Docling)：无法删除临时文件: {file_path}, 原因: {str(del_e)}")
        
        return jsonify({
            'error': error_msg,
            'details': traceback.format_exc()
        }), 500

# API：使用Docling转换为HTML
@app.route('/api/convert-to-html-docling', methods=['POST'])
@swag_from('swagger_docs/convert_to_html_docling.yml')
def api_convert_html_docling():
    start_time = time.time()
    
    if 'file' not in request.files:
        logger.warning("API调用(Docling HTML)：没有文件上传")
        return jsonify({'error': '没有文件上传'}), 400
    
    file = request.files['file']
    
    if file.filename == '':
        logger.warning("API调用(Docling HTML)：未选择文件")
        return jsonify({'error': '未选择文件'}), 400
    
    # 获取文件扩展名
    filename = file.filename
    file_ext = os.path.splitext(filename)[1].lower()
    
    logger.info(f"API调用(Docling HTML)：接收到文件: {filename}, 类型: {file_ext}")
    
    # 获取docling转换器
    docling_converter = converter_factory.get_converter('docling')
    if not hasattr(docling_converter, 'is_available') or not docling_converter.is_available:
        logger.error("API调用(Docling HTML)：Docling转换器不可用")
        return jsonify({'error': 'Docling转换器不可用，请确认已安装docling库'}), 500
    
    # 检查文件扩展名是否支持
    if not docling_converter.is_format_supported(file_ext):
        error_msg = f"Docling不支持的文件类型: {file_ext}"
        logger.warning(f"API调用(Docling HTML)：{error_msg}")
        return jsonify({'error': error_msg}), 400
    
    # 保存上传的文件
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    try:
        file.save(file_path)
        logger.info(f"API调用(Docling HTML)：文件保存成功: {file_path}")
        
        # 检查文件大小
        file_size = os.path.getsize(file_path)
        logger.info(f"API调用(Docling HTML)：文件大小: {file_size / 1024:.2f} KB")
        
        # 使用Docling转换为HTML
        logger.info(f"API调用(Docling HTML)：开始转换文件: {filename}")
        html_content = docling_converter.convert_to_html(file_path)
        
        processing_time = time.time() - start_time
        logger.info(f"API调用(Docling HTML)：文件转换完成，耗时: {processing_time:.2f}秒")
        
        # 删除临时文件
        try:
            os.remove(file_path)
            logger.info(f"API调用(Docling HTML)：临时文件已删除: {file_path}")
        except Exception as e:
            logger.warning(f"API调用(Docling HTML)：无法删除临时文件: {file_path}, 原因: {str(e)}")
            try:
                import gc
                gc.collect()
                os.remove(file_path)
                logger.info(f"API调用(Docling HTML)：临时文件已删除(第二次尝试): {file_path}")
            except Exception as e2:
                logger.error(f"API调用(Docling HTML)：无法删除临时文件(第二次尝试): {file_path}, 原因: {str(e2)}")
        
        # 返回API响应
        return jsonify({
            'html': html_content,
            'filename': filename,
            'file_size': file_size,
            'processing_time': round(processing_time, 2),
            'converter': 'docling'
        })
    except Exception as e:
        error_msg = str(e)
        logger.error(f"API调用(Docling HTML)：转换失败: {error_msg}")
        logger.error(traceback.format_exc())
        
        # 出错时删除临时文件
        if os.path.exists(file_path):
            try:
                os.remove(file_path)
                logger.info(f"API调用(Docling HTML)：清理临时文件: {file_path}")
            except Exception as del_e:
                logger.warning(f"API调用(Docling HTML)：无法删除临时文件: {file_path}, 原因: {str(del_e)}")
        
        return jsonify({
            'error': error_msg,
            'details': traceback.format_exc()
        }), 500

# API：将文件转换为文本并保存到指定目录
@app.route('/api/convert-file', methods=['POST'])
@swag_from('swagger_docs/convert_file.yml')
def api_convert_file():
    start_time = time.time()
    
    if 'file' not in request.files:
        logger.warning("API调用(保存文本)：没有文件上传")
        return jsonify({'error': '没有文件上传'}), 400
    
    if 'output_dir' not in request.form:
        logger.warning("API调用(保存文本)：未指定输出目录")
        return jsonify({'error': '未指定输出目录'}), 400
    
    file = request.files['file']
    output_dir = request.form['output_dir']
    
    if file.filename == '':
        logger.warning("API调用(保存文本)：未选择文件")
        return jsonify({'error': '未选择文件'}), 400
    
    # 检查输出目录是否存在，不存在则创建
    if not os.path.exists(output_dir):
        try:
            os.makedirs(output_dir)
            logger.info(f"API调用(保存文本)：创建输出目录: {output_dir}")
        except Exception as e:
            error_msg = f"无法创建输出目录: {str(e)}"
            logger.error(f"API调用(保存文本)：{error_msg}")
        return jsonify({'error': error_msg}), 500
                
                # 获取文件扩展名
    filename = file.filename
    file_ext = os.path.splitext(filename)[1].lower()
    
    logger.info(f"API调用(保存文本)：接收到文件: {filename}, 类型: {file_ext}")
    
    # 保存上传的文件
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    try:
        file.save(file_path)
        logger.info(f"API调用(保存文本)：文件保存成功: {file_path}")
        
        # 检查文件大小
        file_size = os.path.getsize(file_path)
        logger.info(f"API调用(保存文本)：文件大小: {file_size / 1024:.2f} KB")
        
        # 根据文件类型调用相应的转换函数
        if file_ext in ['.doc', '.docx']:
            logger.info(f"API调用(保存文本)：开始转换Word文档: {filename}")
            text = convert_docx(file_path)
        elif file_ext in ['.xls', '.xlsx']:
            logger.info(f"API调用(保存文本)：开始转换Excel文件: {filename}")
            text = convert_xlsx(file_path)
        elif file_ext in ['.ppt', '.pptx']:
            logger.info(f"API调用(保存文本)：开始转换PowerPoint文件: {filename}")
            text = convert_pptx(file_path)
        elif file_ext == '.pdf':
            logger.info(f"API调用(保存文本)：开始转换PDF文件: {filename}")
            text = convert_pdf(file_path)
        elif file_ext == '.txt':
            logger.info(f"API调用(保存文本)：开始读取文本文件: {filename}")
            text = convert_txt(file_path)
        elif file_ext == '.md':
            logger.info(f"API调用(保存文本)：开始转换Markdown文件: {filename}")
            text = convert_md(file_path)
        elif file_ext in ['.mp3', '.wav']:
            logger.info(f"API调用(保存文本)：开始转换音频文件: {filename}")
            text = "音频文件支持Markdown格式导出，请使用转MD API"
        else:
            error_msg = f"不支持的文件类型: {file_ext}"
            logger.warning(f"API调用(保存文本)：{error_msg}")
            
            # 删除临时文件
            try:
                os.remove(file_path)
                logger.info(f"API调用(保存文本)：临时文件已删除: {file_path}")
            except Exception as e:
                logger.warning(f"API调用(保存文本)：无法删除临时文件: {file_path}, 原因: {str(e)}")
            
            return jsonify({'error': error_msg}), 400
        
        # 保存转换后的文本到输出目录
        output_filename = os.path.splitext(filename)[0] + '.txt'
        output_path = os.path.join(output_dir, output_filename)
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(text)
        
        logger.info(f"API调用(保存文本)：文本已保存到: {output_path}")
        
        processing_time = time.time() - start_time
        logger.info(f"API调用(保存文本)：文件转换完成，耗时: {processing_time:.2f}秒")
        
        # 删除临时文件
        try:
            os.remove(file_path)
            logger.info(f"API调用(保存文本)：临时文件已删除: {file_path}")
        except Exception as e:
            logger.warning(f"API调用(保存文本)：无法删除临时文件: {file_path}, 原因: {str(e)}")
            # 尝试强制关闭文件句柄后再删除
            try:
                import gc
                gc.collect()  # 触发垃圾回收，释放未关闭的文件句柄
                os.remove(file_path)
                logger.info(f"API调用(保存文本)：临时文件已删除(第二次尝试): {file_path}")
            except Exception as e2:
                logger.error(f"API调用(保存文本)：无法删除临时文件(第二次尝试): {file_path}, 原因: {str(e2)}")
                # 不再尝试删除，避免阻塞程序
        
        # 返回API响应
        return jsonify({
            'output_path': output_path,
            'filename': filename,
            'file_size': file_size,
            'processing_time': round(processing_time, 2)
        })
    except Exception as e:
        error_msg = str(e)
        logger.error(f"API调用(保存文本)：转换失败: {error_msg}")
        logger.error(traceback.format_exc())
        
        # 出错时删除临时文件
        if os.path.exists(file_path):
            try:
                os.remove(file_path)
                logger.info(f"API调用(保存文本)：清理临时文件: {file_path}")
            except Exception as del_e:
                logger.warning(f"API调用(保存文本)：无法删除临时文件: {file_path}, 原因: {str(del_e)}")
        
        return jsonify({
            'error': error_msg,
            'details': traceback.format_exc()
        }), 500

# API：将文件转换为Markdown并保存到指定目录
@app.route('/api/convert-to-md-file', methods=['POST'])
@swag_from('swagger_docs/convert_to_md_file.yml')
def api_convert_md_file():
    start_time = time.time()
    
    if 'file' not in request.files:
        logger.warning("API调用(保存MD)：没有文件上传")
        return jsonify({'error': '没有文件上传'}), 400
    
    if 'output_dir' not in request.form:
        logger.warning("API调用(保存MD)：未指定输出目录")
        return jsonify({'error': '未指定输出目录'}), 400
    
    file = request.files['file']
    output_dir = request.form['output_dir']
    
    if file.filename == '':
        logger.warning("API调用(保存MD)：未选择文件")
        return jsonify({'error': '未选择文件'}), 400
    
    # 检查输出目录是否存在，不存在则创建
    if not os.path.exists(output_dir):
        try:
            os.makedirs(output_dir)
            logger.info(f"API调用(保存MD)：创建输出目录: {output_dir}")
        except Exception as e:
            error_msg = f"无法创建输出目录: {str(e)}"
            logger.error(f"API调用(保存MD)：{error_msg}")
        return jsonify({'error': error_msg}), 500
                
                # 获取文件扩展名
    filename = file.filename
    file_ext = os.path.splitext(filename)[1].lower()
    
    logger.info(f"API调用(保存MD)：接收到文件: {filename}, 类型: {file_ext}")
    
    # 检查文件扩展名
    supported_extensions = ['.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx', '.pdf', '.txt', '.md', '.mp3', '.wav']
    if file_ext not in supported_extensions:
        logger.warning(f"API调用(保存MD)：不支持的文件格式: {file_ext}")
        return jsonify({'error': f'不支持的文件格式: {file_ext}'}), 400
    
    # 保存上传的文件
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    try:
        file.save(file_path)
        logger.info(f"API调用(保存MD)：文件保存成功: {file_path}")
        
        # 检查文件大小
        file_size = os.path.getsize(file_path)
        logger.info(f"API调用(保存MD)：文件大小: {file_size / 1024:.2f} KB")
        
        # 使用MarkItDown转换为Markdown
        logger.info(f"API调用(保存MD)：开始将文件转换为Markdown: {filename}")
        markdown_text = convert_to_markdown(file_path)
        
        # 保存转换后的Markdown到输出目录
        output_filename = os.path.splitext(filename)[0] + '.md'
        output_path = os.path.join(output_dir, output_filename)
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(markdown_text)
        
        logger.info(f"API调用(保存MD)：Markdown已保存到: {output_path}")
        
        processing_time = time.time() - start_time
        logger.info(f"API调用(保存MD)：文件转换完成，耗时: {processing_time:.2f}秒")
        
        # 删除临时文件
        try:
            os.remove(file_path)
            logger.info(f"API调用(保存MD)：临时文件已删除: {file_path}")
        except Exception as e:
            logger.warning(f"API调用(保存MD)：无法删除临时文件: {file_path}, 原因: {str(e)}")
            try:
                import gc
                gc.collect()
                os.remove(file_path)
                logger.info(f"API调用(保存MD)：临时文件已删除(第二次尝试): {file_path}")
            except Exception as e2:
                logger.error(f"API调用(保存MD)：无法删除临时文件(第二次尝试): {file_path}, 原因: {str(e2)}")
        
        # 返回API响应
        return jsonify({
            'output_path': output_path,
            'filename': filename,
            'file_size': file_size,
            'processing_time': round(processing_time, 2)
        })
    except Exception as e:
        error_msg = str(e)
        logger.error(f"API调用(保存MD)：转换失败: {error_msg}")
        logger.error(traceback.format_exc())
        
        # 出错时删除临时文件
        if os.path.exists(file_path):
            try:
                os.remove(file_path)
                logger.info(f"API调用(保存MD)：清理临时文件: {file_path}")
            except Exception as del_e:
                logger.warning(f"API调用(保存MD)：无法删除临时文件: {file_path}, 原因: {str(del_e)}")
        
        return jsonify({
            'error': error_msg,
            'details': traceback.format_exc()
        }), 500

# API：使用Docling将文件转换为Markdown并保存到指定目录
@app.route('/api/convert-to-md-file-docling', methods=['POST'])
@swag_from('swagger_docs/convert_to_md_file_docling.yml')
def api_convert_md_file_docling():
    start_time = time.time()
    
    if 'file' not in request.files:
        logger.warning("API调用(Docling保存)：没有文件上传")
        return jsonify({'error': '没有文件上传'}), 400
    
    if 'output_dir' not in request.form:
        logger.warning("API调用(Docling保存)：未指定输出目录")
        return jsonify({'error': '未指定输出目录'}), 400
    
    file = request.files['file']
    output_dir = request.form['output_dir']
    
    if file.filename == '':
        logger.warning("API调用(Docling保存)：未选择文件")
        return jsonify({'error': '未选择文件'}), 400
    
    # 检查输出目录是否存在，不存在则创建
    if not os.path.exists(output_dir):
        try:
            os.makedirs(output_dir)
            logger.info(f"API调用(Docling保存)：创建输出目录: {output_dir}")
        except Exception as e:
            error_msg = f"无法创建输出目录: {str(e)}"
            logger.error(f"API调用(Docling保存)：{error_msg}")
            return jsonify({'error': error_msg}), 500
    
    # 获取文件扩展名
    filename = file.filename
    file_ext = os.path.splitext(filename)[1].lower()
    
    logger.info(f"API调用(Docling保存)：接收到文件: {filename}, 类型: {file_ext}")
    
    # 获取docling转换器
    docling_converter = converter_factory.get_converter('docling')
    if not hasattr(docling_converter, 'is_available') or not docling_converter.is_available:
        logger.error("API调用(Docling保存)：Docling转换器不可用")
        return jsonify({'error': 'Docling转换器不可用，请确认已安装docling库'}), 500
    
    # 检查文件扩展名是否支持
    if not docling_converter.is_format_supported(file_ext):
        error_msg = f"Docling不支持的文件类型: {file_ext}"
        logger.warning(f"API调用(Docling保存)：{error_msg}")
        return jsonify({'error': error_msg}), 400
    
    # 保存上传的文件
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    try:
        file.save(file_path)
        logger.info(f"API调用(Docling保存)：文件保存成功: {file_path}")
        
        # 检查文件大小
        file_size = os.path.getsize(file_path)
        logger.info(f"API调用(Docling保存)：文件大小: {file_size / 1024:.2f} KB")
        
        # 使用Docling转换为Markdown
        logger.info(f"API调用(Docling保存)：开始转换文件: {filename}")
        markdown_text = docling_converter.convert_to_markdown(file_path)
        
        # 保存转换后的Markdown到输出目录
        output_filename = os.path.splitext(filename)[0] + '.md'
        output_path = os.path.join(output_dir, output_filename)
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(markdown_text)
        
        logger.info(f"API调用(Docling保存)：Markdown已保存到: {output_path}")
        
        processing_time = time.time() - start_time
        logger.info(f"API调用(Docling保存)：文件转换完成，耗时: {processing_time:.2f}秒")
        
        # 删除临时文件
        try:
            os.remove(file_path)
            logger.info(f"API调用(Docling保存)：临时文件已删除: {file_path}")
        except Exception as e:
            logger.warning(f"API调用(Docling保存)：无法删除临时文件: {file_path}, 原因: {str(e)}")
            try:
                import gc
                gc.collect()
                os.remove(file_path)
                logger.info(f"API调用(Docling保存)：临时文件已删除(第二次尝试): {file_path}")
            except Exception as e2:
                logger.error(f"API调用(Docling保存)：无法删除临时文件(第二次尝试): {file_path}, 原因: {str(e2)}")
        
        # 返回API响应
        return jsonify({
            'output_path': output_path,
            'filename': filename,
            'file_size': file_size,
            'processing_time': round(processing_time, 2),
            'converter': 'docling'
        })
    except Exception as e:
        error_msg = str(e)
        logger.error(f"API调用(Docling保存)：转换失败: {error_msg}")
        logger.error(traceback.format_exc())
        
        # 出错时删除临时文件
        if os.path.exists(file_path):
            try:
                os.remove(file_path)
                logger.info(f"API调用(Docling保存)：清理临时文件: {file_path}")
            except Exception as del_e:
                logger.warning(f"API调用(Docling保存)：无法删除临时文件: {file_path}, 原因: {str(del_e)}")
        
        return jsonify({
            'error': error_msg,
            'details': traceback.format_exc()
        }), 500

# 工具函数：从文档中提取图片
def extract_images_from_document(file_path, output_dir, base_filename):
    """
    从文档中提取所有图片并保存到指定目录
    
    Args:
        file_path (str): 文档文件路径
        output_dir (str): 图片保存目录
        base_filename (str): 原始文件名，用于构造图片名称
        
    Returns:
        tuple: (图片数量, 图片路径列表)
    """
    image_paths = []
    file_ext = os.path.splitext(file_path)[1].lower()
    timestamp = int(time.time())
    
    # 确保输出目录存在
    os.makedirs(output_dir, exist_ok=True)
    
    # 支持的文件类型
    supported_extensions = ['.docx', '.xls', '.xlsx', '.pptx']
    if file_ext not in supported_extensions:
        raise ValueError(f"不支持的文件类型: {file_ext}, 支持的类型: {', '.join(supported_extensions)}")
    
    try:
        # 从Word文档提取图片
        if file_ext == '.docx':
            from docx import Document
            from docx.parts.image import ImagePart
            
            doc = Document(file_path)
            
            # 获取所有图片附件
            image_parts = []
            
            # 获取文档主体中的所有图片
            for rel_id, rel in doc.part.rels.items():
                if isinstance(rel.target_part, ImagePart):
                    # 收集图片信息，包括rel_id作为标识
                    image_parts.append((rel_id, rel.target_part))
            
            # 扫描文档，记录图片出现的顺序
            ordered_images = []
            
            # 遍历文档中的段落
            for para in doc.paragraphs:
                # 检查段落中的图片
                for run in para.runs:
                    if run._element.xpath('.//a:blip'):
                        for blip in run._element.xpath('.//a:blip'):
                            embed = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                            if embed and embed not in [img[0] for img in ordered_images]:
                                if embed in doc.part.rels:
                                    rel = doc.part.rels[embed]
                                    if isinstance(rel.target_part, ImagePart):
                                        ordered_images.append((embed, rel.target_part))
            
            # 如果通过扫描没有找到所有图片，使用原始收集的图片集合
            if not ordered_images:
                ordered_images = image_parts
            
            # 保存图片，按顺序命名
            for i, (rel_id, img_part) in enumerate(ordered_images, 1):
                image_filename = f"{base_filename}_{timestamp}_photo{i}.png"
                image_path = os.path.join(output_dir, image_filename)
                
                with open(image_path, 'wb') as f:
                    f.write(img_part.blob)
                
                image_paths.append(image_path)
                logger.info(f"提取图片: {image_filename}")
        
        # 从Excel提取图片
        elif file_ext in ['.xls', '.xlsx']:
            import openpyxl
            
            workbook = openpyxl.load_workbook(file_path)
            image_index = 0
            
            # 遍历所有工作表并提取图片
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                
                if hasattr(sheet, '_images'):
                    for image in sheet._images:
                        image_index += 1
                        image_filename = f"{base_filename}_{timestamp}_photo{image_index}.png"
                        image_path = os.path.join(output_dir, image_filename)
                        
                        # 提取并保存图片
                        with open(image_path, 'wb') as f:
                            f.write(image._data())
                        
                        image_paths.append(image_path)
                        logger.info(f"提取图片: {image_filename}")
        
        # 从PowerPoint提取图片
        elif file_ext == '.pptx':
            from pptx import Presentation
            
            prs = Presentation(file_path)
            image_index = 0
            image_dict = {}  # 用于跟踪已保存图片，避免重复
            
            # 遍历所有幻灯片并提取图片
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'image') and shape.image:
                        # 检查是否已保存过此图片（避免重复）
                        image_hash = hash(shape.image.blob)
                        if image_hash not in image_dict:
                            image_index += 1
                            image_filename = f"{base_filename}_{timestamp}_photo{image_index}.png"
                            image_path = os.path.join(output_dir, image_filename)
                            
                            # 提取并保存图片
                            with open(image_path, 'wb') as f:
                                f.write(shape.image.blob)
                            
                            image_paths.append(image_path)
                            image_dict[image_hash] = image_path
                            logger.info(f"提取图片: {image_filename}")
                    
                    # 检查组合图形
                    if hasattr(shape, 'shapes'):
                        for subshape in shape.shapes:
                            if hasattr(subshape, 'image') and subshape.image:
                                image_hash = hash(subshape.image.blob)
                                if image_hash not in image_dict:
                                    image_index += 1
                                    image_filename = f"{base_filename}_{timestamp}_photo{image_index}.png"
                                    image_path = os.path.join(output_dir, image_filename)
                                    
                                    # 提取并保存图片
                                    with open(image_path, 'wb') as f:
                                        f.write(subshape.image.blob)
                                    
                                    image_paths.append(image_path)
                                    image_dict[image_hash] = image_path
                                    logger.info(f"提取图片: {image_filename}")
        
        return len(image_paths), image_paths
    
    except Exception as e:
        raise Exception(f"提取图片过程中出错: {str(e)}")

# 转换为Markdown并导出图片（Docling）
@app.route('/convert-to-md-images-file-docling', methods=['POST'])
def convert_to_md_images_file_docling():
    """
    使用Docling将文件转换为Markdown并导出图片，同时将图片放入静态目录中以支持预览
    """
    start_time = time.time()
    execution_id = datetime.now().strftime("%Y%m%d%H%M%S")
     # 获取当前请求的主机和协议
    host_url = request.host_url.rstrip('/') 
    
    # 获取上传的文件
    if 'file' not in request.files:
        logger.warning(f"Docling图片导出{execution_id}：没有文件上传")
        return jsonify({'error': '没有文件上传'}), 400
    
    file = request.files['file']
    
    if file.filename == '':
        logger.warning(f"Docling图片导出{execution_id}：未选择文件")
        return jsonify({'error': '未选择文件'}), 400
    
    # 获取文件扩展名
    filename = file.filename
    base_filename = os.path.splitext(filename)[0]
    file_ext = os.path.splitext(filename)[1].lower()
    
    logger.info(f"Docling图片导出{execution_id}：接收到文件: {filename}, 类型: {file_ext}")
    
    # 初始化Docling转换器
    docling_converter = converter_factory.get_converter("docling")
    
    # 检查文件扩展名是否支持
    if not docling_converter.is_format_supported(file_ext):
        error_msg = f"Docling不支持的文件类型: {file_ext}"
        logger.warning(f"Docling图片导出{execution_id}：{error_msg}")
        return jsonify({'error': error_msg}), 400
    
    # 保存上传的文件
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    
    # 创建保存图片的静态目录
    static_img_dir = os.path.join('app', 'static', 'images', 'exported', execution_id)
    static_img_url_path = f"{host_url}/static/images/exported/{execution_id}"
    
    try:
        file.save(file_path)
        logger.info(f"Docling图片导出{execution_id}：文件保存成功: {file_path}")
        
        # 检查文件大小
        file_size = os.path.getsize(file_path)
        logger.info(f"Docling图片导出{execution_id}：文件大小: {file_size / 1024:.2f} KB")
        
        # 创建静态图片目录
        if not os.path.exists(static_img_dir):
            os.makedirs(static_img_dir, exist_ok=True)
            logger.info(f"Docling图片导出{execution_id}：创建静态图片目录: {static_img_dir}")
        
        # 使用Docling处理文件
        try:
            from docling.document_converter import DocumentConverter
            from docling.datamodel.pipeline_options import PdfPipelineOptions
            from docling.datamodel.base_models import InputFormat, FigureElement
            from docling.document_converter import PdfFormatOption
            from docling_core.types.doc import ImageRefMode, PictureItem, TableItem
            
            # 设置Docling选项，启用图片导出
            IMAGE_RESOLUTION_SCALE = 2.0
            pipeline_options = PdfPipelineOptions()
            pipeline_options.images_scale = IMAGE_RESOLUTION_SCALE
            pipeline_options.generate_page_images = True
            pipeline_options.generate_picture_images = True
            
            # 创建通用格式选项
            format_options = {}
            
            # 根据文件类型添加适当的格式选项
            if file_ext.lower() in ['.pdf']:
                format_options[InputFormat.PDF] = PdfFormatOption(pipeline_options=pipeline_options)
            # 对于其他格式，不设置特定选项，使用Docling默认设置
            
            doc_converter = DocumentConverter(format_options=format_options)
            
            logger.info(f"Docling图片导出{execution_id}：开始转换文件并导出图片")
            doc_converter = DocumentConverter(
                format_options={
                    InputFormat.PDF: PdfFormatOption(pipeline_options=pipeline_options)
                }
            )
            
            logger.info(f"Docling图片导出{execution_id}：开始转换文件并导出图片")
            
            # 转换文件
            conv_res = doc_converter.convert(file_path)
            doc_filename = conv_res.input.file.stem
            
            # 保存Markdown文件
            md_filename_full = os.path.join(static_img_dir, f"{doc_filename}.md")
            
            # 先使用独立引用模式保存原始的markdown
            conv_res.document.save_as_markdown(md_filename_full, image_mode=ImageRefMode.REFERENCED)
            logger.info(f"Docling图片导出{execution_id}：Markdown已保存到: {md_filename_full}")
            
            # 计数器
            page_counter = 0
            table_counter = 0
            picture_counter = 0
            
            # 保存页面图片
            page_image_paths = []
            page_image_urls = []
            for page_no, page in conv_res.document.pages.items():
                page_no = page.page_no
                page_image_filename = f"{doc_filename}-page-{page_no}.png"
                page_image_path = os.path.join(static_img_dir, page_image_filename)
                page_image_url = f"{static_img_url_path}/{page_image_filename}"
                
                try:
                    if hasattr(page, 'image') and page.image is not None and hasattr(page.image, 'pil_image') and page.image.pil_image is not None:
                        with open(page_image_path, 'wb') as fp:
                            page.image.pil_image.save(fp, format="PNG")
                        page_image_paths.append(page_image_path)
                        page_image_urls.append(page_image_url)
                        page_counter += 1
                        logger.info(f"Docling图片导出{execution_id}：页面图片已保存: {page_image_path}")
                    else:
                        logger.warning(f"Docling图片导出{execution_id}：页面 {page_no} 没有可用的图片")
                except Exception as e:
                    logger.warning(f"Docling图片导出{execution_id}：保存页面 {page_no} 图片时出错: {str(e)}")
            
            # 保存图片和表格
            table_image_paths = []
            table_image_urls = []
            picture_image_paths = []
            picture_image_urls = []
            
            for element, _level in conv_res.document.iterate_items():
                if isinstance(element, TableItem):
                    try:
                        image = element.get_image(conv_res.document)
                        if image is not None:
                            table_counter += 1  # 只有当成功获取图片时才增加计数
                            table_image_filename = f"{doc_filename}-table-{table_counter}.png"
                            table_image_path = os.path.join(static_img_dir, table_image_filename)
                            table_image_url = f"{static_img_url_path}/{table_image_filename}"
                            
                            with open(table_image_path, 'wb') as fp:
                                image.save(fp, "PNG")
                            table_image_paths.append(table_image_path)
                            table_image_urls.append(table_image_url)
                            logger.info(f"Docling图片导出{execution_id}：表格图片已保存: {table_image_path}")
                        else:
                            logger.warning(f"Docling图片导出{execution_id}：无法获取表格图片，图片为None")
                    except Exception as e:
                        logger.warning(f"Docling图片导出{execution_id}：处理表格图片时出错: {str(e)}")
                
                if isinstance(element, PictureItem):
                    try:
                        image = element.get_image(conv_res.document)
                        if image is not None:
                            picture_counter += 1  # 只有当成功获取图片时才增加计数
                            picture_image_filename = f"{doc_filename}-picture-{picture_counter}.png"
                            picture_image_path = os.path.join(static_img_dir, picture_image_filename)
                            picture_image_url = f"{static_img_url_path}/{picture_image_filename}"
                            
                            with open(picture_image_path, 'wb') as fp:
                                image.save(fp, "PNG")
                            picture_image_paths.append(picture_image_path)
                            picture_image_urls.append(picture_image_url)
                            logger.info(f"Docling图片导出{execution_id}：图片已保存: {picture_image_path}")
                        else:
                            logger.warning(f"Docling图片导出{execution_id}：无法获取图片，图片为None")
                    except Exception as e:
                        logger.warning(f"Docling图片导出{execution_id}：处理图片时出错: {str(e)}")
            
            # 读取原始Markdown内容
            with open(md_filename_full, 'r', encoding='utf-8') as f:
                markdown_text = f.read()
            
            # 将相对路径替换为完整URL路径
            # 提取artifacts文件夹名
            artifacts_dir = f"{doc_filename}_artifacts"
            
            # 替换 ![Image](artifacts_dir/imagename.png) 模式的图片路径
            pattern = r'!\[([^\]]*)\]\((' + re.escape(artifacts_dir) + r'[^)]+)\)'
            replacement = r'![\1](' + static_img_url_path + r'/\2)'
            
            markdown_text = re.sub(pattern, replacement, markdown_text)

            # 修复带有 backslashes 的 Windows 风格路径
            def replace_with_forward_slash(match):
                path = match.group(1)
                return f'![Image]({path.replace("\\", "/")})'
            
            # 匹配图片链接中的路径并替换反斜杠
            markdown_text = re.sub(r'!\[Image\]\(([^)]+)\)', replace_with_forward_slash, markdown_text)
            
            # 替换页面图片路径
            for i, path in enumerate(page_image_paths):
                markdown_text = markdown_text.replace(path, page_image_urls[i])
            
            # 替换表格图片路径
            for i, path in enumerate(table_image_paths):
                markdown_text = markdown_text.replace(path, table_image_urls[i])
            
            # 替换普通图片路径
            for i, path in enumerate(picture_image_paths):
                markdown_text = markdown_text.replace(path, picture_image_urls[i])
            
            # 保存修改后的Markdown内容
            with open(md_filename_full, 'w', encoding='utf-8') as f:
                f.write(markdown_text)
            
            processing_time = time.time() - start_time
            logger.info(f"Docling图片导出{execution_id}：处理完成，耗时: {processing_time:.2f}秒")
            
            # 返回处理结果
            return jsonify({
                'text': markdown_text,
                'processing_time': round(processing_time, 2),
                'converter': 'docling',
                'page_count': page_counter,
                'table_count': table_counter,
                'picture_count': picture_counter
            })
            
        except ImportError as e:
            error_msg = f"Docling库导入错误: {str(e)}"
            logger.error(f"Docling图片导出{execution_id}：{error_msg}")
            return jsonify({'error': error_msg}), 500
        except Exception as e:
            error_msg = f"Docling处理失败: {str(e)}"
            logger.error(f"Docling图片导出{execution_id}：{error_msg}")
            logger.error(traceback.format_exc())
            return jsonify({'error': error_msg, 'details': traceback.format_exc()}), 500
    
    except Exception as e:
        error_msg = f"处理文件时出错: {str(e)}"
        logger.error(f"Docling图片导出{execution_id}：{error_msg}")
        logger.error(traceback.format_exc())
        
        # 删除临时文件
        if os.path.exists(file_path):
            try:
                os.remove(file_path)
                logger.info(f"Docling图片导出{execution_id}：临时文件已删除: {file_path}")
            except Exception as del_e:
                logger.warning(f"Docling图片导出{execution_id}：无法删除临时文件: {file_path}, 原因: {str(del_e)}")
        
        return jsonify({'error': error_msg, 'details': traceback.format_exc()}), 500
    
    finally:
        # 删除临时文件
        if os.path.exists(file_path):
            try:
                os.remove(file_path)
                logger.info(f"Docling图片导出{execution_id}：临时文件已删除: {file_path}")
            except Exception as del_e:
                logger.warning(f"Docling图片导出{execution_id}：无法删除临时文件: {file_path}, 原因: {str(del_e)}")

@app.route('/about')
def about_page():
    logger.info("访问关于页面")
    return render_template('about.html')

# API：将URL转换为Markdown
@app.route('/api/convert-url-to-md', methods=['POST'])
@swag_from('swagger_docs/convert_url_to_md.yml')
def api_convert_url_to_md():
    """
    将URL网页内容转换为Markdown格式
    """
    start_time = time.time()
    
    # 检查是否提供了URL
    if 'url' not in request.form:
        logger.warning("API调用(URL转MD)：未提供URL")
        return jsonify({'error': '未提供URL'}), 400
    
    url = request.form['url']
    remove_header_footer = request.form.get('remove_header_footer', 'true').lower() in ['true', '1', 't', 'y', 'yes']
    
    logger.info(f"API调用(URL转MD)：开始处理URL: {url}，移除页眉页脚: {remove_header_footer}")
    
    # 初始化URL转换器
    url_converter = URLConverter()
    
    try:
        # 将URL转换为Markdown
        markdown_text = url_converter.convert_url_to_markdown(url, remove_header_footer)
        
        processing_time = time.time() - start_time
        logger.info(f"API调用(URL转MD)：URL转换完成，耗时: {processing_time:.2f}秒")
        
        # 返回API响应
        return jsonify({
            'text': markdown_text,
            'url': url,
            'processing_time': round(processing_time, 2)
        })
    except Exception as e:
        error_msg = str(e)
        logger.error(f"API调用(URL转MD)：转换失败: {error_msg}")
        logger.error(traceback.format_exc())
        
        return jsonify({
            'error': error_msg,
            'details': traceback.format_exc()
        }), 500

# API：将URL转换为Markdown文件并保存
@app.route('/api/convert-url-to-md-file', methods=['POST'])
@swag_from('swagger_docs/convert_url_to_md_file.yml')
def api_convert_url_to_md_file():
    """
    将URL网页内容转换为Markdown格式并保存为文件
    """
    start_time = time.time()
    
    # 检查是否提供了URL和输出目录
    if 'url' not in request.form:
        logger.warning("API调用(URL转MD文件)：未提供URL")
        return jsonify({'error': '未提供URL'}), 400
    
    if 'output_dir' not in request.form:
        logger.warning("API调用(URL转MD文件)：未提供输出目录")
        return jsonify({'error': '未提供输出目录'}), 400
    
    url = request.form['url']
    output_dir = request.form['output_dir']
    remove_header_footer = request.form.get('remove_header_footer', 'true').lower() in ['true', '1', 't', 'y', 'yes']
    
    logger.info(f"API调用(URL转MD文件)：开始处理URL: {url}，输出目录: {output_dir}，移除页眉页脚: {remove_header_footer}")
    
    # 确保输出目录存在
    try:
        os.makedirs(output_dir, exist_ok=True)
        logger.info(f"确保输出目录存在: {output_dir}")
    except Exception as e:
        error_msg = f"创建输出目录失败: {str(e)}"
        logger.error(error_msg)
        return jsonify({'error': error_msg}), 500
    
    # 初始化URL转换器
    url_converter = URLConverter()
    
    try:
        # 将URL转换为Markdown
        markdown_text = url_converter.convert_url_to_markdown(url, remove_header_footer)
        
        # 从URL中提取文件名
        import re
        from urllib.parse import urlparse
        
        parsed_url = urlparse(url)
        domain = parsed_url.netloc
        path = parsed_url.path
        
        # 生成安全的文件名
        if path and path != '/':
            # 从路径中提取最后一部分作为文件名
            filename_base = path.strip('/').split('/')[-1]
            # 清理文件名
            filename_base = re.sub(r'[\\/*?:"<>|]', '_', filename_base)
            if not filename_base:
                filename_base = domain
        else:
            filename_base = domain
        
        # 如果文件名太长或为空，则使用域名
        if len(filename_base) > 100 or not filename_base:
            filename_base = domain.replace('.', '_')
        
        # 添加时间戳避免重名
        timestamp = datetime.now().strftime("%Y%m%d%H%M%S")
        filename = f"{filename_base}_{timestamp}.md"
        
        # 保存为Markdown文件
        output_path = os.path.join(output_dir, filename)
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(markdown_text)
        
        processing_time = time.time() - start_time
        logger.info(f"API调用(URL转MD文件)：URL转换完成并保存到: {output_path}，耗时: {processing_time:.2f}秒")
        
        # 返回API响应
        return jsonify({
            'output_path': output_path,
            'url': url,
            'filename': filename,
            'processing_time': round(processing_time, 2)
        })
    except Exception as e:
        error_msg = str(e)
        logger.error(f"API调用(URL转MD文件)：转换失败: {error_msg}")
        logger.error(traceback.format_exc())
        
        return jsonify({
            'error': error_msg,
            'details': traceback.format_exc()
        }), 500

if __name__ == '__main__':
    logger.info("应用启动")
    app.run(host='0.0.0.0', port=5000,debug=True)
    logger.info("应用关闭") 