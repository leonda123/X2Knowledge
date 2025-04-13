"""
通用工具函数模块
"""
import os
import logging
import time
import traceback
from datetime import datetime

from app import logger

def extract_images_from_document(file_path, output_dir, base_filename):
    """
    从文档中提取所有图片并保存到指定目录
    
    @param {str} file_path - 文档文件路径
    @param {str} output_dir - 图片保存目录
    @param {str} base_filename - 原始文件名，用于构造图片名称
    @returns {tuple} - (图片数量, 图片路径列表)
    """
    image_paths = []
    file_ext = os.path.splitext(file_path)[1].lower()
    timestamp = int(time.time())
    
    # 确保输出目录存在
    os.makedirs(output_dir, exist_ok=True)
    
    # 支持的文件类型
    supported_extensions = ['.docx', '.xls', '.xlsx', '.pptx']
    if file_ext not in supported_extensions:
        raise ValueError(f"不支持的文件类型: {file_ext}, 支持的类型: {', '.join(supported_extensions)}")
    
    try:
        # 从Word文档提取图片
        if file_ext == '.docx':
            from docx import Document
            from docx.parts.image import ImagePart
            
            doc = Document(file_path)
            
            # 获取所有图片附件
            image_parts = []
            
            # 获取文档主体中的所有图片
            for rel_id, rel in doc.part.rels.items():
                if isinstance(rel.target_part, ImagePart):
                    # 收集图片信息，包括rel_id作为标识
                    image_parts.append((rel_id, rel.target_part))
            
            # 扫描文档，记录图片出现的顺序
            ordered_images = []
            
            # 遍历文档中的段落
            for para in doc.paragraphs:
                # 检查段落中的图片
                for run in para.runs:
                    if run._element.xpath('.//a:blip'):
                        for blip in run._element.xpath('.//a:blip'):
                            embed = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                            if embed and embed not in [img[0] for img in ordered_images]:
                                if embed in doc.part.rels:
                                    rel = doc.part.rels[embed]
                                    if isinstance(rel.target_part, ImagePart):
                                        ordered_images.append((embed, rel.target_part))
            
            # 如果通过扫描没有找到所有图片，使用原始收集的图片集合
            if not ordered_images:
                ordered_images = image_parts
            
            # 保存图片，按顺序命名
            for i, (rel_id, img_part) in enumerate(ordered_images, 1):
                image_filename = f"{base_filename}_{timestamp}_photo{i}.png"
                image_path = os.path.join(output_dir, image_filename)
                
                with open(image_path, 'wb') as f:
                    f.write(img_part.blob)
                
                image_paths.append(image_path)
                logger.info(f"提取图片: {image_filename}")
        
        # 从Excel提取图片
        elif file_ext in ['.xls', '.xlsx']:
            import openpyxl
            
            workbook = openpyxl.load_workbook(file_path)
            image_index = 0
            
            # 遍历所有工作表并提取图片
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                
                if hasattr(sheet, '_images'):
                    for image in sheet._images:
                        image_index += 1
                        image_filename = f"{base_filename}_{timestamp}_photo{image_index}.png"
                        image_path = os.path.join(output_dir, image_filename)
                        
                        # 提取并保存图片
                        with open(image_path, 'wb') as f:
                            f.write(image._data())
                        
                        image_paths.append(image_path)
                        logger.info(f"提取图片: {image_filename}")
        
        # 从PowerPoint提取图片
        elif file_ext == '.pptx':
            from pptx import Presentation
            
            prs = Presentation(file_path)
            image_index = 0
            image_dict = {}  # 用于跟踪已保存图片，避免重复
            
            # 遍历所有幻灯片并提取图片
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'image') and shape.image:
                        # 检查是否已保存过此图片（避免重复）
                        image_hash = hash(shape.image.blob)
                        if image_hash not in image_dict:
                            image_index += 1
                            image_filename = f"{base_filename}_{timestamp}_photo{image_index}.png"
                            image_path = os.path.join(output_dir, image_filename)
                            
                            # 提取并保存图片
                            with open(image_path, 'wb') as f:
                                f.write(shape.image.blob)
                            
                            image_paths.append(image_path)
                            image_dict[image_hash] = image_path
                            logger.info(f"提取图片: {image_filename}")
                    
                    # 检查组合图形
                    if hasattr(shape, 'shapes'):
                        for subshape in shape.shapes:
                            if hasattr(subshape, 'image') and subshape.image:
                                image_hash = hash(subshape.image.blob)
                                if image_hash not in image_dict:
                                    image_index += 1
                                    image_filename = f"{base_filename}_{timestamp}_photo{image_index}.png"
                                    image_path = os.path.join(output_dir, image_filename)
                                    
                                    # 提取并保存图片
                                    with open(image_path, 'wb') as f:
                                        f.write(subshape.image.blob)
                                    
                                    image_paths.append(image_path)
                                    image_dict[image_hash] = image_path
                                    logger.info(f"提取图片: {image_filename}")
        
        return len(image_paths), image_paths
    
    except Exception as e:
        raise Exception(f"提取图片过程中出错: {str(e)}")

def cleanup_temp_file(file_path, context=""):
    """
    清理临时文件，带有重试机制
    
    @param {str} file_path - 要删除的文件路径
    @param {str} context - 上下文信息，用于日志
    """
    if os.path.exists(file_path):
        try:
            os.remove(file_path)
            logger.info(f"{context}临时文件已删除: {file_path}")
        except Exception as e:
            logger.warning(f"{context}无法删除临时文件: {file_path}, 原因: {str(e)}")
            # 尝试强制关闭文件句柄后再删除
            try:
                import gc
                gc.collect()  # 触发垃圾回收，释放未关闭的文件句柄
                os.remove(file_path)
                logger.info(f"{context}临时文件已删除(第二次尝试): {file_path}")
            except Exception as e2:
                logger.error(f"{context}无法删除临时文件(第二次尝试): {file_path}, 原因: {str(e2)}")
                # 不再尝试删除，避免阻塞程序 