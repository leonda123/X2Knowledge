import docx
import openpyxl
from pptx import Presentation
import PyPDF2
import markdown
import os
import logging
import traceback
import sys
import subprocess
import tempfile
import io
from PIL import Image

# 导入OCR工具
try:
    from app.utils.ocr_utils import extract_text_from_image, extract_text_from_image_bytes, extract_text_from_pdf_images, HAS_TESSERACT
except ImportError:
    HAS_TESSERACT = False
    logger = logging.getLogger(__name__)
    logger.warning("无法导入OCR工具模块，图片文字识别功能将不可用")

# 导入MarkItDown
try:
    from markitdown import MarkItDown
    HAS_MARKITDOWN = True
    logger = logging.getLogger(__name__)
    logger.info("成功导入MarkItDown库")
except ImportError:
    HAS_MARKITDOWN = False
    logger = logging.getLogger(__name__)
    logger.warning("无法导入MarkItDown库，将使用替代方法进行文档转Markdown")

# 获取日志记录器
logger = logging.getLogger(__name__)

# 尝试导入处理.doc文件所需的库
HAS_TEXTRACT = False
HAS_WIN32 = False

try:
    import win32com.client
    HAS_WIN32 = sys.platform == 'win32'
    logger.info("成功导入win32com库")
except ImportError:
    logger.warning("无法导入win32com库，将无法使用COM方式处理.doc文件")

try:
    import textract
    HAS_TEXTRACT = True
    logger.info("成功导入textract库")
except ImportError:
    logger.warning("无法导入textract库，将使用替代方法处理.doc文件")

# 尝试导入音频处理库
try:
    import pydub
    HAS_PYDUB = True
    logger.info("成功导入pydub库")
except ImportError:
    HAS_PYDUB = False
    logger.warning("无法导入pydub库，将无法处理音频文件")

def convert_docx(file_path):
    """将Word文档转换为文本"""
    try:
        logger.info(f"开始处理Word文档: {file_path}")
        
        # 检查文件是否存在
        if not os.path.exists(file_path):
            error_msg = f"文件不存在: {file_path}"
            logger.error(error_msg)
            raise FileNotFoundError(error_msg)
        
        # 获取文件扩展名
        file_ext = os.path.splitext(file_path)[1].lower()
        
        # 处理.doc格式（旧版Word）
        if file_ext == '.doc':
            logger.info("检测到旧版.doc格式，使用替代方法处理")
            return convert_doc_to_text(file_path)
        
        # 处理.docx格式（新版Word）
        logger.info("处理.docx格式文件")
        try:
            doc = docx.Document(file_path)
            full_text = []
            
            # 记录段落数量
            paragraph_count = len(doc.paragraphs)
            logger.info(f"文档包含 {paragraph_count} 个段落")
            
            # 提取段落文本
            for i, para in enumerate(doc.paragraphs):
                if i % 100 == 0 and i > 0:
                    logger.debug(f"已处理 {i}/{paragraph_count} 个段落")
                
                # 确保段落文本使用UTF-8编码
                para_text = para.text
                if para_text:
                    # 检查是否包含非ASCII字符
                    if any(ord(c) > 127 for c in para_text):
                        logger.debug(f"段落 {i} 包含非ASCII字符，确保UTF-8编码")
                    full_text.append(para_text)
                else:
                    full_text.append("")
            
            # 提取表格内容
            table_count = len(doc.tables)
            logger.info(f"文档包含 {table_count} 个表格")
            
            if table_count > 0:
                full_text.append("\n--- 表格内容 ---\n")
                
                for t_idx, table in enumerate(doc.tables):
                    full_text.append(f"表格 #{t_idx+1}:")
                    
                    # 获取表格行列数
                    row_count = len(table.rows)
                    col_count = len(table.columns)
                    logger.info(f"表格 #{t_idx+1} 包含 {row_count} 行, {col_count} 列")
                    
                    # 提取表格内容
                    for r_idx, row in enumerate(table.rows):
                        row_texts = []
                        for c_idx, cell in enumerate(row.cells):
                            # 确保单元格文本使用UTF-8编码
                            cell_text = cell.text.strip()
                            if cell_text and any(ord(c) > 127 for c in cell_text):
                                logger.debug(f"表格 #{t_idx+1} 单元格 ({r_idx},{c_idx}) 包含非ASCII字符，确保UTF-8编码")
                            row_texts.append(cell_text)
                        
                        # 使用制表符分隔单元格内容
                        full_text.append(" | ".join(row_texts))
                    
                    full_text.append("")  # 表格之间添加空行
            
            # 处理文档中的图片
            if HAS_TESSERACT:
                logger.info("开始处理文档中的图片")
                image_texts = extract_images_from_docx(doc, file_path)
                if image_texts:
                    full_text.append("\n--- 图片中的文本 ---\n")
                    for idx, img_text in enumerate(image_texts):
                        full_text.append(f"图片 #{idx+1}:")
                        # 确保OCR识别的文本使用UTF-8编码
                        full_text.append(img_text)
                        full_text.append("")
            else:
                logger.warning("Tesseract OCR未安装，跳过图片文字识别")
            
            result = '\n'.join(full_text)
            logger.info(f"Word文档处理完成，提取了 {len(result)} 个字符")
            return result
        except Exception as e:
            logger.error(f"处理.docx文件时出错: {str(e)}")
            logger.error(traceback.format_exc())
            return f"[无法读取文档内容: {str(e)}]"
    except Exception as e:
        error_msg = f"Word文档转换失败: {str(e)}"
        logger.error(error_msg)
        logger.error(traceback.format_exc())
        raise Exception(error_msg)

def extract_images_from_docx(doc, file_path):
    """从Word文档中提取图片并识别文本"""
    try:
        logger.info("开始从Word文档提取图片")
        image_texts = []
        
        # 创建临时目录存储图片
        with tempfile.TemporaryDirectory() as temp_dir:
            # 提取文档中的所有关系
            image_count = 0
            
            # 遍历文档中的所有部分
            for rel in doc.part.rels.values():
                # 检查是否是图片
                if "image" in rel.target_ref:
                    try:
                        image_count += 1
                        logger.info(f"处理图片 #{image_count}: {rel.target_ref}")
                        
                        # 获取图片数据
                        image_part = rel.target_part
                        image_bytes = image_part.blob
                        
                        # 保存图片到临时文件
                        image_path = os.path.join(temp_dir, f"image_{image_count}.png")
                        with open(image_path, "wb") as f:
                            f.write(image_bytes)
                        
                        # 提取图片文本
                        text = extract_text_from_image(image_path)
                        if text and text != "[图片中未检测到文本]":
                            image_texts.append(text)
                    except Exception as e:
                        logger.warning(f"处理图片 #{image_count} 时出错: {str(e)}")
            
            logger.info(f"从Word文档中提取了 {image_count} 张图片，识别出 {len(image_texts)} 张图片中的文本")
        
        return image_texts
    except Exception as e:
        logger.error(f"提取Word文档图片时出错: {str(e)}")
        logger.error(traceback.format_exc())
        return []

def convert_doc_to_text(file_path):
    """使用替代方法将.doc文件转换为文本"""
    word_app = None
    doc = None
    try:
        logger.info(f"尝试转换.doc文件: {file_path}")
        
        # 方法1: 使用textract库
        if HAS_TEXTRACT:
            try:
                logger.info("尝试使用textract库提取文本")
                # 尝试检测编码
                text_bytes = textract.process(file_path)
                
                # 尝试不同的编码解码文本
                encodings = ['utf-8', 'gbk', 'gb2312', 'gb18030', 'big5', 'cp936']
                for encoding in encodings:
                    try:
                        text = text_bytes.decode(encoding)
                        # 检查是否有中文字符
                        if any('\u4e00' <= char <= '\u9fff' for char in text):
                            logger.info(f"使用 {encoding} 成功解码中文内容")
                            logger.info(f"使用textract成功提取文本，共 {len(text)} 个字符")
                            return text
                    except UnicodeDecodeError:
                        continue
                
                # 如果没有找到合适的编码，使用默认的utf-8 with errors='replace'
                text = text_bytes.decode('utf-8', errors='replace')
                logger.info(f"使用textract成功提取文本，共 {len(text)} 个字符")
                return text
            except Exception as e:
                logger.warning(f"textract提取失败: {str(e)}，尝试其他方法")
        else:
            logger.info("textract库不可用，跳过此方法")
        
        # 方法2: 在Windows上使用pywin32
        if HAS_WIN32:
            try:
                logger.info("尝试使用pywin32提取文本")
                # 创建临时文件
                temp_dir = tempfile.mkdtemp()
                temp_docx = os.path.join(temp_dir, "temp.docx")
                
                # 使用Word转换doc到docx
                word_app = win32com.client.Dispatch("Word.Application")
                word_app.visible = False
                
                # 设置Word应用程序的编码
                try:
                    # 尝试设置默认编码为简体中文
                    word_app.DefaultSaveEncoding = 936  # 936是简体中文GBK的代码页
                except:
                    logger.warning("无法设置Word默认编码")
                
                doc = word_app.Documents.Open(os.path.abspath(file_path))
                
                # 保存为docx格式，确保保留表格和格式
                try:
                    # 使用SaveAs2方法，保留所有格式
                    doc.SaveAs2(os.path.abspath(temp_docx), 16)  # 16表示docx格式
                    logger.info(f"成功将.doc文件转换为.docx格式: {temp_docx}")
                except Exception as save_e:
                    logger.warning(f"使用SaveAs2方法保存失败: {str(save_e)}，尝试使用SaveAs方法")
                    doc.SaveAs(os.path.abspath(temp_docx), 16)
                    logger.info(f"使用SaveAs方法成功将.doc文件转换为.docx格式: {temp_docx}")
                
                doc.Close()
                doc = None  # 清除引用
                word_app.Quit()
                word_app = None  # 清除引用
                
                # 读取转换后的docx
                text = convert_docx(temp_docx)
                
                # 清理临时文件
                try:
                    os.remove(temp_docx)
                    os.rmdir(temp_dir)
                except Exception as clean_e:
                    logger.warning(f"清理临时文件失败: {str(clean_e)}")
                
                logger.info(f"使用pywin32成功提取文本，共 {len(text)} 个字符")
                return text
            except Exception as e:
                logger.warning(f"pywin32提取失败: {str(e)}，尝试其他方法")
                # 确保关闭Word应用程序
                if doc is not None:
                    try:
                        doc.Close(False)
                    except:
                        pass
                if word_app is not None:
                    try:
                        word_app.Quit()
                    except:
                        pass
        else:
            logger.info("pywin32不可用或非Windows系统，跳过此方法")
        
        # 方法3: 使用外部工具antiword (Linux/Mac)
        if sys.platform != 'win32':
            try:
                logger.info("尝试使用antiword提取文本")
                text = subprocess.check_output(['antiword', '-m', 'UTF-8', file_path]).decode('utf-8', errors='replace')
                logger.info(f"使用antiword成功提取文本，共 {len(text)} 个字符")
                return text
            except Exception as e:
                logger.warning(f"antiword提取失败: {str(e)}")
        
        # 方法4: 使用catdoc (Linux/Mac)
        if sys.platform != 'win32':
            try:
                logger.info("尝试使用catdoc提取文本")
                text = subprocess.check_output(['catdoc', '-d', 'utf-8', file_path]).decode('utf-8', errors='replace')
                logger.info(f"使用catdoc成功提取文本，共 {len(text)} 个字符")
                return text
            except Exception as e:
                logger.warning(f"catdoc提取失败: {str(e)}")
        
        # 方法5: 尝试直接读取文件内容
        try:
            logger.info("尝试直接读取文件内容")
            with open(file_path, 'rb') as f:
                content = f.read()
                # 尝试提取可读文本
                text = ""
                # 中文编码尝试列表
                encodings = ['utf-8', 'gbk', 'gb2312', 'gb18030', 'big5', 'cp936', 'latin-1', 'cp1252']
                
                for encoding in encodings:
                    try:
                        logger.debug(f"尝试使用 {encoding} 编码解码文件内容")
                        text = content.decode(encoding, errors='replace')
                        # 检查是否有中文字符
                        if any('\u4e00' <= char <= '\u9fff' for char in text):
                            logger.info(f"使用 {encoding} 编码成功解码中文内容")
                            break
                    except Exception as e:
                        logger.debug(f"{encoding} 编码解码失败: {str(e)}")
                        continue
                
                # 过滤掉不可打印字符，但保留中文字符
                filtered_text = ''
                for c in text:
                    if c.isprintable() or c in ['\n', '\t', '\r'] or '\u4e00' <= c <= '\u9fff':
                        filtered_text += c
                
                logger.info(f"直接读取提取了 {len(filtered_text)} 个字符")
                if len(filtered_text.strip()) > 0:
                    return filtered_text
        except Exception as e:
            logger.warning(f"直接读取失败: {str(e)}")
        
        # 如果所有方法都失败，返回错误信息
        error_msg = "无法提取.doc文件内容，请将文件转换为.docx格式后重试"
        logger.error(error_msg)
        return f"[{error_msg}]"
        
    except Exception as e:
        error_msg = f".doc文件转换失败: {str(e)}"
        logger.error(error_msg)
        logger.error(traceback.format_exc())
        return f"[{error_msg}]"
    finally:
        # 确保关闭Word应用程序和文档
        if doc is not None:
            try:
                doc.Close(False)
                logger.debug("Word文档已关闭")
            except:
                pass
        if word_app is not None:
            try:
                word_app.Quit()
                logger.debug("Word应用程序已关闭")
            except:
                pass

def convert_xlsx(file_path):
    """将Excel文件转换为文本"""
    wb = None
    try:
        logger.info(f"开始处理Excel文件: {file_path}")
        
        # 检查文件是否存在
        if not os.path.exists(file_path):
            error_msg = f"文件不存在: {file_path}"
            logger.error(error_msg)
            raise FileNotFoundError(error_msg)
            
        wb = openpyxl.load_workbook(file_path, data_only=True)
        text_content = []
        
        # 记录工作表数量
        sheet_count = len(wb.sheetnames)
        logger.info(f"Excel文件包含 {sheet_count} 个工作表")
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            
            # 确保工作表名称使用UTF-8编码
            if any(ord(c) > 127 for c in sheet_name):
                logger.debug(f"工作表名称 '{sheet_name}' 包含非ASCII字符，确保UTF-8编码")
            
            text_content.append(f"工作表: {sheet_name}")
            logger.debug(f"处理工作表: {sheet_name}")
            
            # 获取工作表的行列数
            row_count = sheet.max_row
            col_count = sheet.max_column
            logger.debug(f"工作表 {sheet_name} 包含 {row_count} 行, {col_count} 列")
            
            for row in sheet.iter_rows(values_only=True):
                # 处理每个单元格，确保使用UTF-8编码
                processed_cells = []
                for cell in row:
                    cell_text = str(cell) if cell is not None else ''
                    
                    # 检查是否包含非ASCII字符
                    if cell_text and any(ord(c) > 127 for c in cell_text):
                        logger.debug(f"单元格内容包含非ASCII字符，确保UTF-8编码")
                    
                    processed_cells.append(cell_text)
                
                row_text = '\t'.join(processed_cells)
                text_content.append(row_text)
            
            text_content.append('\n')
        
        result = '\n'.join(text_content)
        logger.info(f"Excel文件处理完成，提取了 {len(result)} 个字符")
        return result
    except Exception as e:
        error_msg = f"Excel文件转换失败: {str(e)}"
        logger.error(error_msg)
        logger.error(traceback.format_exc())
        raise Exception(error_msg)
    finally:
        # 确保关闭工作簿，释放文件句柄
        if wb is not None:
            try:
                wb._archive.close()
                logger.debug("Excel工作簿已关闭")
            except Exception as e:
                logger.warning(f"关闭Excel工作簿时出错: {str(e)}")

def convert_pptx(file_path):
    """将PowerPoint文件转换为文本"""
    try:
        logger.info(f"开始处理PowerPoint文件: {file_path}")
        
        # 检查文件是否存在
        if not os.path.exists(file_path):
            error_msg = f"文件不存在: {file_path}"
            logger.error(error_msg)
            raise FileNotFoundError(error_msg)
            
        prs = Presentation(file_path)
        text_content = []
        
        # 记录幻灯片数量
        slide_count = len(prs.slides)
        logger.info(f"PowerPoint文件包含 {slide_count} 张幻灯片")
        
        # 创建临时目录存储图片
        with tempfile.TemporaryDirectory() as temp_dir:
            # 处理每张幻灯片
            for i, slide in enumerate(prs.slides):
                text_content.append(f"幻灯片 #{i+1}")
                logger.debug(f"处理幻灯片 #{i+1}")
                
                shape_count = 0
                text_shape_count = 0
                image_count = 0
                
                # 提取形状中的文本
                for shape in slide.shapes:
                    shape_count += 1
                    
                    # 提取文本形状
                    if hasattr(shape, "text") and shape.text:
                        text_shape_count += 1
                        
                        # 确保形状文本使用UTF-8编码
                        shape_text = shape.text
                        if any(ord(c) > 127 for c in shape_text):
                            logger.debug(f"幻灯片 #{i+1} 形状文本包含非ASCII字符，确保UTF-8编码")
                        
                        text_content.append(shape_text)
                    
                    # 提取图片
                    if HAS_TESSERACT and shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                        try:
                            image_count += 1
                            logger.info(f"处理幻灯片 #{i+1} 中的图片 #{image_count}")
                            
                            # 保存图片到临时文件
                            image_path = os.path.join(temp_dir, f"slide_{i+1}_image_{image_count}.png")
                            
                            # 获取图片数据
                            image = shape.image
                            image_bytes = image.blob
                            
                            with open(image_path, "wb") as f:
                                f.write(image_bytes)
                            
                            # 提取图片文本
                            img_text = extract_text_from_image(image_path)
                            if img_text and img_text != "[图片中未检测到文本]":
                                text_content.append(f"[图片 #{image_count} 文本:]")
                                
                                # 确保OCR识别的文本使用UTF-8编码
                                if any(ord(c) > 127 for c in img_text):
                                    logger.debug(f"幻灯片 #{i+1} 图片 #{image_count} 文本包含非ASCII字符，确保UTF-8编码")
                                
                                text_content.append(img_text)
                        except Exception as e:
                            logger.warning(f"处理幻灯片图片时出错: {str(e)}")
                
                logger.debug(f"幻灯片 #{i+1} 包含 {shape_count} 个形状，其中 {text_shape_count} 个包含文本，{image_count} 个图片")
                text_content.append('\n')
        
        result = '\n'.join(text_content)
        logger.info(f"PowerPoint文件处理完成，提取了 {len(result)} 个字符")
        return result
    except Exception as e:
        error_msg = f"PowerPoint文件转换失败: {str(e)}"
        logger.error(error_msg)
        logger.error(traceback.format_exc())
        raise Exception(error_msg)

def convert_pdf(file_path):
    """将PDF文件转换为文本"""
    pdf_file = None
    try:
        logger.info(f"开始处理PDF文件: {file_path}")
        
        # 检查文件是否存在
        if not os.path.exists(file_path):
            error_msg = f"文件不存在: {file_path}"
            logger.error(error_msg)
            raise FileNotFoundError(error_msg)
            
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            pdf_file = file  # 保存文件引用以便在finally中关闭
            text_content = []
            
            # 记录页面数量
            page_count = len(reader.pages)
            logger.info(f"PDF文件包含 {page_count} 页")
            
            # 提取文本内容
            for i in range(page_count):
                page = reader.pages[i]
                logger.debug(f"处理第 {i+1} 页")
                
                try:
                    page_text = page.extract_text()
                    
                    # 确保页面文本使用UTF-8编码
                    if page_text:
                        # 检查是否包含非ASCII字符
                        if any(ord(c) > 127 for c in page_text):
                            logger.debug(f"页面 {i+1} 包含非ASCII字符，确保UTF-8编码")
                        
                        text_content.append(f"页面 #{i+1}")
                        text_content.append(page_text)
                    else:
                        text_content.append(f"页面 #{i+1}")
                        text_content.append("[此页无文本内容]")
                except Exception as page_error:
                    logger.warning(f"提取第 {i+1} 页文本时出错: {str(page_error)}")
                    text_content.append(f"页面 #{i+1} [提取文本失败]")
                
                text_content.append('\n')
            
            # 处理PDF中的图片（如果有）
            if HAS_TESSERACT:
                logger.info("开始处理PDF中的图片")
                image_texts = extract_text_from_pdf_images(file_path)
                if image_texts:
                    text_content.append("--- PDF图片中的文本 ---")
                    for img_data in image_texts:
                        text_content.append(f"页面 #{img_data['page']} 图片文本:")
                        
                        # 确保OCR识别的文本使用UTF-8编码
                        img_text = img_data['text']
                        if any(ord(c) > 127 for c in img_text):
                            logger.debug(f"页面 {img_data['page']} 图片文本包含非ASCII字符，确保UTF-8编码")
                        
                        text_content.append(img_text)
                        text_content.append("")
            else:
                logger.warning("Tesseract OCR未安装，跳过PDF图片文字识别")
            
            result = '\n'.join(text_content)
            logger.info(f"PDF文件处理完成，提取了 {len(result)} 个字符")
            return result
    except Exception as e:
        error_msg = f"PDF文件转换失败: {str(e)}"
        logger.error(error_msg)
        logger.error(traceback.format_exc())
        raise Exception(error_msg)
    finally:
        # 确保关闭PDF文件
        if pdf_file is not None and not pdf_file.closed:
            try:
                pdf_file.close()
                logger.debug("PDF文件已关闭")
            except Exception as e:
                logger.warning(f"关闭PDF文件时出错: {str(e)}")

def convert_txt(file_path):
    """读取文本文件内容"""
    try:
        logger.info(f"开始处理文本文件: {file_path}")
        
        # 检查文件是否存在
        if not os.path.exists(file_path):
            error_msg = f"文件不存在: {file_path}"
            logger.error(error_msg)
            raise FileNotFoundError(error_msg)
            
        # 尝试不同的编码方式
        encodings = ['utf-8', 'gbk', 'gb2312', 'latin-1']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as file:
                    content = file.read()
                    logger.info(f"文本文件成功以 {encoding} 编码读取，共 {len(content)} 个字符")
                    return content
            except UnicodeDecodeError:
                logger.debug(f"尝试以 {encoding} 编码读取失败，尝试下一种编码")
                continue
        
        # 如果所有编码都失败，使用二进制模式读取并使用errors='replace'
        logger.warning("所有编码尝试失败，使用替换模式读取")
        with open(file_path, 'r', encoding='utf-8', errors='replace') as file:
            content = file.read()
            logger.info(f"文本文件以替换模式读取，共 {len(content)} 个字符")
            return content
    except Exception as e:
        error_msg = f"文本文件读取失败: {str(e)}"
        logger.error(error_msg)
        logger.error(traceback.format_exc())
        raise Exception(error_msg)

def convert_md(file_path):
    """将Markdown文件转换为纯文本"""
    try:
        logger.info(f"开始处理Markdown文件: {file_path}")
        
        # 检查文件是否存在
        if not os.path.exists(file_path):
            error_msg = f"文件不存在: {file_path}"
            logger.error(error_msg)
            raise FileNotFoundError(error_msg)
            
        # 尝试不同的编码方式
        encodings = ['utf-8', 'gbk', 'latin-1']
        
        md_content = None
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as file:
                    md_content = file.read()
                    logger.debug(f"Markdown文件成功以 {encoding} 编码读取")
                    break
            except UnicodeDecodeError:
                logger.debug(f"尝试以 {encoding} 编码读取失败，尝试下一种编码")
                continue
        
        if md_content is None:
            # 如果所有编码都失败，使用替换模式
            logger.warning("所有编码尝试失败，使用替换模式读取")
            with open(file_path, 'r', encoding='utf-8', errors='replace') as file:
                md_content = file.read()
        
        logger.info(f"Markdown文件读取完成，开始转换为HTML")
        # 将Markdown转换为HTML
        html = markdown.markdown(md_content)
        
        logger.debug("HTML生成完成，开始转换为纯文本")
        # 简单替换一些HTML标签为文本等效形式
        text = html.replace('<h1>', '# ').replace('</h1>', '\n')
        text = text.replace('<h2>', '## ').replace('</h2>', '\n')
        text = text.replace('<h3>', '### ').replace('</h3>', '\n')
        text = text.replace('<p>', '').replace('</p>', '\n')
        text = text.replace('<strong>', '').replace('</strong>', '')
        text = text.replace('<em>', '').replace('</em>', '')
        text = text.replace('<code>', '`').replace('</code>', '`')
        text = text.replace('<pre>', '```\n').replace('</pre>', '\n```')
        text = text.replace('<br>', '\n').replace('<br/>', '\n')
        text = text.replace('<ul>', '').replace('</ul>', '')
        text = text.replace('<li>', '- ').replace('</li>', '')
        
        logger.info(f"Markdown文件处理完成，提取了 {len(text)} 个字符")
        return text
    except Exception as e:
        error_msg = f"Markdown文件转换失败: {str(e)}"
        logger.error(error_msg)
        logger.error(traceback.format_exc())
        raise Exception(error_msg)

# 添加MarkItDown转换功能
def convert_to_markdown(file_path):
    """使用MarkItDown将文件转换为Markdown格式"""
    try:
        logger.info(f"使用MarkItDown开始处理文件: {file_path}")
        
        # 检查文件是否存在
        if not os.path.exists(file_path):
            error_msg = f"文件不存在: {file_path}"
            logger.error(error_msg)
            raise FileNotFoundError(error_msg)
        
        if HAS_MARKITDOWN:
            # 初始化MarkItDown
            md = MarkItDown(enable_plugins=False)
            
            # 转换文件
            result = md.convert(file_path)
            
            # 获取Markdown内容
            markdown_content = result.text_content
            
            logger.info(f"MarkItDown转换成功，生成了 {len(markdown_content)} 个字符的Markdown内容")
            return markdown_content
        else:
            logger.warning("MarkItDown库不可用，使用替代方法")
            # 根据文件类型使用替代方法
            file_ext = os.path.splitext(file_path)[1].lower()
            
            if file_ext in ['.doc', '.docx']:
                text = convert_docx(file_path)
                return f"# {os.path.basename(file_path)}\n\n{text}"
            elif file_ext in ['.xls', '.xlsx']:
                text = convert_xlsx(file_path)
                return f"# {os.path.basename(file_path)}\n\n{text}"
            elif file_ext in ['.ppt', '.pptx']:
                text = convert_pptx(file_path)
                return f"# {os.path.basename(file_path)}\n\n{text}"
            elif file_ext == '.pdf':
                text = convert_pdf(file_path)
                return f"# {os.path.basename(file_path)}\n\n{text}"
            elif file_ext == '.txt':
                text = convert_txt(file_path)
                return f"# {os.path.basename(file_path)}\n\n{text}"
            elif file_ext == '.md':
                return convert_md(file_path)
            elif file_ext in ['.mp3', '.wav']:
                if HAS_PYDUB:
                    return convert_audio_to_markdown(file_path)
                else:
                    return f"# {os.path.basename(file_path)}\n\n[音频文件，无法转换：需要安装pydub库]"
            else:
                return f"# {os.path.basename(file_path)}\n\n[不支持的文件格式: {file_ext}]"
    except Exception as e:
        error_msg = f"MarkItDown转换失败: {str(e)}"
        logger.error(error_msg)
        logger.error(traceback.format_exc())
        raise Exception(error_msg)

def convert_audio_to_markdown(file_path):
    """将音频文件转换为Markdown格式描述"""
    try:
        logger.info(f"开始处理音频文件: {file_path}")
        
        # 获取文件基本信息
        file_name = os.path.basename(file_path)
        file_size = os.path.getsize(file_path)
        file_size_mb = file_size / (1024 * 1024)
        
        # 使用pydub获取音频详细信息
        if HAS_PYDUB:
            try:
                from pydub import AudioSegment
                audio = AudioSegment.from_file(file_path)
                
                # 获取音频时长（毫秒转换为分:秒格式）
                duration_ms = len(audio)
                duration_sec = duration_ms / 1000
                minutes = int(duration_sec // 60)
                seconds = int(duration_sec % 60)
                
                # 获取声道数和采样率
                channels = audio.channels
                frame_rate = audio.frame_rate
                
                # 构建Markdown内容
                md_content = f"# 音频文件: {file_name}\n\n"
                md_content += f"## 音频信息\n\n"
                md_content += f"- **文件名**: {file_name}\n"
                md_content += f"- **文件大小**: {file_size_mb:.2f} MB\n"
                md_content += f"- **时长**: {minutes}分{seconds}秒\n"
                md_content += f"- **声道数**: {channels}\n"
                md_content += f"- **采样率**: {frame_rate} Hz\n"
                
                logger.info(f"音频文件处理完成: {file_name}, 时长: {minutes}分{seconds}秒")
                return md_content
            except Exception as e:
                logger.error(f"处理音频详细信息时出错: {str(e)}")
                
                # 如果pydub处理失败，返回基本信息
                md_content = f"# 音频文件: {file_name}\n\n"
                md_content += f"## 音频信息\n\n"
                md_content += f"- **文件名**: {file_name}\n"
                md_content += f"- **文件大小**: {file_size_mb:.2f} MB\n"
                md_content += f"- **注意**: 无法获取详细音频信息\n"
                
                return md_content
        else:
            # 如果没有pydub，返回基本信息
            md_content = f"# 音频文件: {file_name}\n\n"
            md_content += f"## 音频信息\n\n"
            md_content += f"- **文件名**: {file_name}\n"
            md_content += f"- **文件大小**: {file_size_mb:.2f} MB\n"
            md_content += f"- **注意**: 需安装pydub库获取详细音频信息\n"
            
            return md_content
    except Exception as e:
        error_msg = f"音频文件转换失败: {str(e)}"
        logger.error(error_msg)
        logger.error(traceback.format_exc())
        return f"# 音频文件转换错误\n\n处理时出现错误: {str(e)}" 